import json
import asyncio
import logging
import os
import re
import base64
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any

import httpx
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, ReplyKeyboardMarkup, KeyboardButton, BotCommand, WebAppInfo
from telegram.ext import Application, MessageHandler, CommandHandler, CallbackQueryHandler, filters, ContextTypes

# =============================================
#  API KEYS
# =============================================
TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
GROQ_KEY = os.getenv("GROQ_API_KEY", "")
GEMINI_KEY = os.getenv("GEMINI_API_KEY", "")
DEEPSEEK_KEY = os.getenv("DEEPSEEK_API_KEY", "")
COHERE_KEY = os.getenv("COHERE_API_KEY", "")
OPENROUTER_KEY = os.getenv("OPENROUTER_API_KEY", "")
SERPER_KEY = os.getenv("SERPER_API_KEY", "")
WEBAPP_URL = os.getenv("WEBAPP_URL", "https://kaionavan.github.io/Tgbot1/")
WEATHER_KEY = os.getenv("WEATHER_API_KEY", "7d2fe9b0de0b0dc5c9db043f5705f1a7")
YOUTUBE_KEY = os.getenv("YOUTUBE_API_KEY", "AQ.Ab8RN6IUxZwllKer_PjKZzma2BpDldoRR0QMDwWNKgg4b5V6Qg")
ELEVENLABS_KEY = os.getenv("ELEVENLABS_API_KEY", "sk_4f36ded5d3261c33de95f257934a1246f2e956d699767a13")
OWNER_CHAT_ID = int(os.getenv("OWNER_CHAT_ID", "0"))  # твой chat_id для утренней рассылки

if not TELEGRAM_TOKEN:
    raise RuntimeError("TELEGRAM_TOKEN не задан!")

# =============================================
#  FILES
# =============================================
DATA_FILE = "data.json"
REMINDERS_FILE = "reminders.json"

# =============================================
#  LOGGING
# =========================================23
logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# =============================================
#  TEMPLATES
# =============================================
TEMPLATES = {
    "🌐 Сайт": "Напиши полный HTML файл с красивым CSS дизайном и JS, всё в одном файле. Сделай: ",
    "🐍 Python": "Напиши полный рабочий Python код без сокращений. Задача: ",
    "🤖 Telegram бот": "Напиши полный код Telegram бота на Python. Бот должен: ",
    "🎮 Игра": "Напиши полную игру на Python или HTML/JS. Игра: ",
    "📱 Приложение": "Напиши полное приложение, весь код целиком. Приложение: ",
    "🗄️ База данных": "Напиши Python код с SQLite базой данных. Функционал: ",
    "🎯 C++": "Напиши полный рабочий C++ код. Задача: ",
    "⚡ JavaScript": "Напиши полный JavaScript код в одном файле. Задача: ",
    "📚 Объяснить": "Объясни простым языком с примерами: ",
    "🔍 Поиск": "Найди актуальную информацию в интернете про: ",
}

# =============================================
#  KEYBOARDS
# =============================================
def main_keyboard() -> ReplyKeyboardMarkup:
    """Главная клавиатура с кнопками"""
    webapp_url = WEBAPP_URL + ("?groq=" + GROQ_KEY if GROQ_KEY else "")
    keyboard = [
        # Главная кнопка — голосовой агент + карты
        [KeyboardButton("🎙 Голос + Карты", web_app=WebAppInfo(url=webapp_url))],
        # Чаты
        [KeyboardButton("💬 Новый чат"), KeyboardButton("📂 Мои чаты"), KeyboardButton("🔍 Поиск")],
        # Быстрые действия
        [KeyboardButton("🌤 Погода"), KeyboardButton("🔗 Ссылка"), KeyboardButton("📄 Файл")],
        # Продуктивность
        [KeyboardButton("✅ Задачи"), KeyboardButton("🎯 Цели"), KeyboardButton("🏃 Привычки")],
        # Знания
        [KeyboardButton("📚 База знаний"), KeyboardButton("🎓 Репетитор"), KeyboardButton("🔬 Исследовать")],
        # Развитие
        [KeyboardButton("📖 Книги"), KeyboardButton("🤝 Дебаты"), KeyboardButton("💼 Собеседование")],
        # AI инструменты
        [KeyboardButton("🤖 Агент"), KeyboardButton("📧 Email"), KeyboardButton("📍 Рядом")],
        # Профиль и настройки
        [KeyboardButton("👤 Профиль"), KeyboardButton("📊 Статистика"), KeyboardButton("❓ Помощь")],
    ]
    return ReplyKeyboardMarkup(
    keyboard,
    resize_keyboard=True,
    is_persistent=True,
    input_field_placeholder="Выбери действие или напиши сообщение..."
    )
# =============================================
#  DATA MANAGEMENT
# =============================================
def load_data() -> Dict[str, Any]:
    """Загружает данные из файла data.json"""
    if Path(DATA_FILE).exists():
        try:
            with open(DATA_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Error loading data: {e}")
    
    # Структура по умолчанию
    return {
        "current": "main",
        "chats": {
            "main": {
                "name": "Основной чат",
                "history": [],
                "created": datetime.now().strftime("%d.%m.%Y")
            }
        },
        "topics": [],
        "last_code": "",
        "profile": {
            "name": "",
            "city": "",
            "interests": [],
            "goals": [],
            "schedule": {},
            "created": datetime.now().strftime("%d.%m.%Y")
        },
        "tasks": [],
        "stats": {
            "messages_total": 0,
            "voice_total": 0,
            "files_total": 0,
            "links_total": 0,
            "topics_studied": 0,
            "tasks_done": 0,
            "days_active": 0,
            "first_seen": datetime.now().strftime("%d.%m.%Y"),
            "last_seen": datetime.now().strftime("%d.%m.%Y"),
            "hourly": {},
            "daily_messages": {}
        },
        "habits": []
    }

def save_data(data: Dict[str, Any]) -> None:
    """Сохраняет данные в файл data.json"""
    try:
        with open(DATA_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Error saving data: {e}")

def load_reminders() -> List[Dict[str, Any]]:
    """Загружает напоминания из файла reminders.json"""
    if Path(REMINDERS_FILE).exists():
        try:
            with open(REMINDERS_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Error loading reminders: {e}")
    return []

def save_reminders(reminders: List[Dict[str, Any]]) -> None:
    """Сохраняет напоминания в файл reminders.json"""
    try:
        with open(REMINDERS_FILE, "w", encoding="utf-8") as f:
            json.dump(reminders, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Error saving reminders: {e}")

def get_current_history(data: Dict[str, Any]) -> List[Dict[str, str]]:
    """Возвращает историю текущего чата"""
    return data["chats"][data["current"]]["history"]

def add_message(data: Dict[str, Any], role: str, text: str) -> None:
    """Добавляет сообщение в историю текущего чата"""
    history = get_current_history(data)
    history.append({
        "role": role,
        "text": text[:500],  # Ограничиваем длину
        "time": datetime.now().strftime("%d.%m %H:%M")
    })
    # Храним только последние 50 сообщений
    if len(history) > 50:
        data["chats"][data["current"]]["history"] = history[-50:]

def build_context(data: Dict[str, Any]) -> str:
    """Строит расширенный контекст для AI с долгосрочной памятью"""
    history = get_current_history(data)[-10:]
    history_text = "\n".join([f"{msg['role']}: {msg['text']}" for msg in history])
    topics = ", ".join(data["topics"][-10:]) if data["topics"] else "пока ничего"
    chat_name = data["chats"][data["current"]]["name"]
    last_code = data.get("last_code", "")
    code_context = f"\nПоследний код:\n{last_code[:400]}" if last_code else ""

    # Долгосрочная память
    memory_ctx = build_memory_context(data)
    memory_block = f"\n\n=== ДОЛГОСРОЧНАЯ ПАМЯТЬ ===\n{memory_ctx}" if memory_ctx else ""

    # Персонализация стиля
    style_block = build_personalized_system(data)
    style_hint = f"\n\n=== СТИЛЬ ПОЛЬЗОВАТЕЛЯ ===\n{style_block}" if style_block else ""

    profile = data.get("profile", {})
    user_name = profile.get("name", "")
    address = f"Обращайся по имени: {user_name}." if user_name else "Обращайся: господин."
    city = profile.get("city") or data.get("weather_city", "Ташкент")
    
    now = datetime.utcnow() + timedelta(hours=data.get("timezone_offset", 5))
    time_ctx = f"Сейчас: {now.strftime('%d.%m.%Y %H:%M')} (UTC+{data.get('timezone_offset',5)}, {city})"
    
    goals = data.get("goals", [])
    active_goals = [g["title"] for g in goals if g.get("status") == "active"][:3]
    goals_ctx = f"Активные цели: {', '.join(active_goals)}." if active_goals else ""
    
    tasks = [t["text"] for t in data.get("tasks", []) if not t.get("done")][:3]
    tasks_ctx = f"Задачи на сегодня: {', '.join(tasks)}." if tasks else ""

    return (
        f"Ты личный AI-агент высокого класса. {address} Язык: русский.\n"
        f"{time_ctx}. Чат: {chat_name}.\n"
        f"Темы пользователя: {topics}.\n"
        f"{goals_ctx} {tasks_ctx}\n"
        f"{memory_block}{style_hint}\n"
        f"История диалога:\n{history_text}{code_context}\n"
        f"ПРАВИЛА:\n"
        f"1. Код — полностью, без сокращений, с комментариями.\n"
        f"2. Ответы краткие и конкретные — без воды.\n"
        f"3. После ответа ВСЕГДА предлагай следующий шаг (что можно сделать дальше).\n"
        f"4. Если видишь что задача сложная — разбей на шаги сам.\n"
        f"5. Обращайся по-человечески, как личный помощник, а не робот.\n"
        f"6. Если пользователь написал задачу/цель — предложи добавить в трекер."
    )

def build_code_prompt(task: str) -> str:
    """Строит промпт для генерации кода"""
    return (
        f"Задача: {task}\n\n"
        f"ТРЕБОВАНИЯ К КОДУ:\n"
        f"1. Напиши ПОЛНЫЙ код от первой до последней строки\n"
        f"2. НЕ используй заглушки типа # TODO или # здесь ваш код\n"
        f"3. Добавь все необходимые импорты\n"
        f"4. Код должен запускаться без изменений\n"
        f"5. Добавь обработку ошибок\n"
        f"6. Комментарии на русском языке\n"
        f"7. Весь код должен быть в одном файле\n"
        f"8. Минимум 50 строк кода\n\n"
        f"Напиши код:"
    )

def extract_code_from_text(text: str) -> Optional[str]:
    """Извлекает код из текста (между ```)"""
    # Ищем блоки кода
    patterns = [
        r'```(?:\w+)?\n(.*?)```',
        r'```(.*?)```',
        r'`(.*?)`',
    ]
    
    for pattern in patterns:
        codes = re.findall(pattern, text, re.DOTALL)
        if codes:
            # Берем самый длинный блок
            best = max(codes, key=len).strip()
            if len(best) > 50:  # Минимальная длина кода
                return best
    
    return None

def detect_extension(prompt: str, code: str) -> str:
    """Определяет расширение файла по промпту и коду"""
    prompt_lower = prompt.lower()
    
    # По промпту
    if any(x in prompt_lower for x in ['html', 'сайт', 'веб-страница', 'веб страница']):
        return 'html'
    if 'css' in prompt_lower:
        return 'css'
    if any(x in prompt_lower for x in ['javascript', 'js', 'скрипт']):
        return 'js'
    if any(x in prompt_lower for x in ['c++', 'cpp', 'плюсы']):
        return 'cpp'
    if any(x in prompt_lower for x in ['c#', 'csharp', 'шарп']):
        return 'cs'
    if 'kotlin' in prompt_lower:
        return 'kt'
    if 'swift' in prompt_lower:
        return 'swift'
    if 'rust' in prompt_lower:
        return 'rs'
    if any(x in prompt_lower for x in ['golang', 'go ']):
        return 'go'
    if 'php' in prompt_lower:
        return 'php'
    if 'ruby' in prompt_lower:
        return 'rb'
    if any(x in prompt_lower for x in ['bash', 'shell', 'sh']):
        return 'sh'
    if 'sql' in prompt_lower:
        return 'sql'
    if any(x in prompt_lower for x in ['dart', 'flutter']):
        return 'dart'
    if 'java' in prompt_lower and 'javascript' not in prompt_lower:
        return 'java'
    
    # По содержимому кода
    if code:
        if '#include' in code or 'cout' in code or 'using namespace' in code:
            return 'cpp'
        if 'using System' in code or 'namespace' in code and 'class' in code:
            return 'cs'
        if 'console.log' in code or 'function' in code and 'var' in code:
            return 'js'
        if '<html' in code or '<!DOCTYPE' in code:
            return 'html'
        if 'def ' in code and 'import' in code:
            return 'py'
        if 'package main' in code and 'func main' in code:
            return 'go'
        if '<?php' in code:
            return 'php'
        if 'CREATE TABLE' in code or 'SELECT * FROM' in code:
            return 'sql'
    
    # По умолчанию
    return 'py'

def is_code_request(text: str) -> bool:
    """Проверяет, просит ли пользователь код"""
    keywords = [
        "создай", "напиши", "сделай", "код", "приложение", "скрипт",
        "программу", "бот", "сайт", "функцию", "калькулятор", "игру",
        "база данных", "класс", "алгоритм", "парсер", "api"
    ]
    text_lower = text.lower()
    return any(k in text_lower for k in keywords)

def is_search_request(text: str) -> bool:
    """Проверяет, просит ли пользователь поиск"""
    keywords = [
        "найди", "поищи", "что сейчас", "новости", "актуально", 
        "курс", "погода", "цена", "кто такой", "что такое", "где"
    ]
    text_lower = text.lower()
    return any(k in text_lower for k in keywords)

# =============================================
#  AI PROVIDERS
# =============================================
async def ask_groq(prompt: str, system_context: str) -> Optional[str]:
    """Запрос к Groq API"""
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            response = await client.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {GROQ_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "llama-3.3-70b-versatile",
                    "messages": [
                        {"role": "system", "content": system_context},
                        {"role": "user", "content": prompt}
                    ],
                    "max_tokens": 8000,
                    "temperature": 0.3
                }
            )
            data = response.json()
            return data["choices"][0]["message"]["content"]
    except Exception as e:
        logger.error(f"Groq error: {e}")
        return None

async def ask_gemini(prompt: str, system_context: str) -> Optional[str]:
    """Запрос к Gemini API"""
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            response = await client.post(
                f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_KEY}",
                json={
                    "contents": [
                        {
                            "parts": [
                                {"text": f"{system_context}\n\n{prompt}\n\nПиши ПОЛНЫЙ код без сокращений!"}
                            ]
                        }
                    ],
                    "generationConfig": {
                        "maxOutputTokens": 8192,
                        "temperature": 0.3
                    }
                }
            )
            data = response.json()
            return data["candidates"][0]["content"]["parts"][0]["text"]
    except Exception as e:
        logger.error(f"Gemini error: {e}")
        return None

async def ask_openrouter(prompt: str, system_context: str) -> Optional[str]:
    """Запрос к OpenRouter API"""
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            response = await client.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {OPENROUTER_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "deepseek/deepseek-chat",
                    "messages": [
                        {"role": "system", "content": system_context},
                        {"role": "user", "content": prompt}
                    ],
                    "max_tokens": 8000
                }
            )
            data = response.json()
            return data["choices"][0]["message"]["content"]
    except Exception as e:
        logger.error(f"OpenRouter error: {e}")
        return None

async def ask_cohere(prompt: str, system_context: str) -> Optional[str]:
    """Запрос к Cohere API"""
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            response = await client.post(
                "https://api.cohere.ai/v1/chat",
                headers={
                    "Authorization": f"Bearer {COHERE_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "command-r-plus",
                    "message": prompt,
                    "preamble": system_context,
                    "max_tokens": 4000
                }
            )
            data = response.json()
            return data["text"]
    except Exception as e:
        logger.error(f"Cohere error: {e}")
        return None

async def ask_deepseek(prompt: str, system_context: str) -> Optional[str]:
    """Запрос к DeepSeek API"""
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            response = await client.post(
                "https://api.deepseek.com/chat/completions",
                headers={
                    "Authorization": f"Bearer {DEEPSEEK_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": "deepseek-chat",
                    "messages": [
                        {"role": "system", "content": system_context},
                        {"role": "user", "content": prompt}
                    ],
                    "max_tokens": 8000
                }
            )
            data = response.json()
            return data["choices"][0]["message"]["content"]
    except Exception as e:
        logger.error(f"DeepSeek error: {e}")
        return None

async def get_best_ai_response(prompt: str, system_context: str, for_code: bool = False) -> Tuple[Optional[str], Optional[str]]:
    """Умный роутер: мультиагент для сложных вопросов, быстрый для простых"""

    # Для сложных вопросов — мультиагентный подход
    if not for_code and is_complex_question(prompt):
        logger.info("Using multi-agent for complex question")
        result, name = await multi_agent_response(prompt, system_context)
        if result:
            return result, name

    # Для простых — по очереди (быстро)
    providers = [
        (ask_groq, "Groq"),
        (ask_gemini, "Gemini"),
        (ask_deepseek, "DeepSeek"),
        (ask_cohere, "Cohere"),
        (ask_openrouter, "OpenRouter"),
    ]
    for provider_func, provider_name in providers:
        try:
            result = await provider_func(prompt, system_context)
            if result:
                logger.info(f"Got response from {provider_name}")
                return result, provider_name
        except Exception as e:
            logger.error(f"{provider_name} error: {e}")
            continue
    return None, None

async def generate_code_smart(task: str, system_context: str) -> Tuple[Optional[str], Optional[str]]:
    """Умная генерация кода с проверкой длины"""
    
    # Первая попытка
    result, ai_name = await get_best_ai_response(
        build_code_prompt(task), 
        system_context, 
        for_code=True
    )
    
    if not result:
        return None, None
    
    # Извлекаем код
    code = extract_code_from_text(result)
    
    # Если код слишком короткий, просим дописать
    if code and len(code) < 800:
        expand_prompt = (
            f"Код слишком короткий. Задача: {task}\n\n"
            f"Предыдущий код:\n```\n{code}\n```\n\n"
            f"Допиши код полностью - добавь все функции, UI, обработку ошибок. "
            f"Напиши весь код заново, но теперь ПОЛНОСТЬЮ, минимум 100 строк!"
        )
        
        result2, ai_name2 = await get_best_ai_response(
            expand_prompt,
            system_context,
            for_code=True
        )
        
        if result2 and len(result2) > len(result):
            return result2, f"{ai_name}+доп"
    
    return result, ai_name

async def check_all_providers() -> Dict[str, str]:
    """Параллельная проверка всех провайдеров"""
    
    test_prompt = "Скажи 'ок' одним словом"
    test_context = "Ты помощник. Отвечай кратко."
    
    providers = [
        ("🟢 Groq", ask_groq),
        ("💎 Gemini", ask_gemini),
        ("🔷 OpenRouter", ask_openrouter),
        ("🟡 Cohere", ask_cohere),
        ("🔵 DeepSeek", ask_deepseek),
    ]
    
    async def check_provider(name: str, func: callable) -> Tuple[str, str]:
        try:
            result = await asyncio.wait_for(func(test_prompt, test_context), timeout=8)
            if result:
                return name, "✅ Работает"
            else:
                return name, "❌ Нет ответа"
        except asyncio.TimeoutError:
            return name, "⏱ Таймаут"
        except Exception:
            return name, "❌ Ошибка"
    
    results = await asyncio.gather(*[
        check_provider(name, func) for name, func in providers
    ])
    
    return dict(results)

# =============================================
#  VOICE AND IMAGE PROCESSING
# =============================================
async def transcribe_voice(audio_data: bytes) -> str:
    """Распознаёт голос через Groq Whisper"""
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                "https://api.groq.com/openai/v1/audio/transcriptions",
                headers={"Authorization": f"Bearer {GROQ_KEY}"},
                files={
                    "file": ("audio.ogg", audio_data, "audio/ogg"),
                    "model": (None, "whisper-large-v3"),
                    "language": (None, "ru"),
                    "response_format": (None, "json")
                }
            )
            data = r.json()
            text = data.get("text", "").strip()
            if text:
                logger.info(f"Voice transcribed: {text[:50]}")
                return text
    except Exception as e:
        logger.error(f"Whisper error: {e}")

    # Fallback — пробуем без указания языка
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                "https://api.groq.com/openai/v1/audio/transcriptions",
                headers={"Authorization": f"Bearer {GROQ_KEY}"},
                files={
                    "file": ("audio.mp3", audio_data, "audio/mpeg"),
                    "model": (None, "whisper-large-v3"),
                }
            )
            return r.json().get("text", "").strip()
    except Exception as e:
        logger.error(f"Whisper fallback error: {e}")
    return ""


async def analyze_image(image_bytes: bytes, prompt: str = "") -> Optional[str]:
    """Анализирует изображение через Groq Vision (llama-4-scout)"""
    import base64
    b64 = base64.b64encode(image_bytes).decode()
    question = prompt or "Опиши подробно что видишь на изображении. Если есть текст — прочитай его весь. Отвечай на русском языке."

    # Способ 1: Groq Vision — llama-4-scout
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers={"Authorization": f"Bearer {GROQ_KEY}", "Content-Type": "application/json"},
                json={
                    "model": "meta-llama/llama-4-scout-17b-16e-instruct",
                    "messages": [{"role": "user", "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
                        {"type": "text", "text": question}
                    ]}],
                    "max_tokens": 1024
                }
            )
            data = r.json()
            if "choices" in data:
                result = data["choices"][0]["message"]["content"]
                if result and len(result) > 5:
                    logger.info("✅ Image analyzed via Groq Vision")
                    return result
            logger.error(f"Groq vision response: {data}")
    except Exception as e:
        logger.error(f"Groq vision error: {e}")

    # Способ 2: Gemini напрямую как резерв
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_KEY}",
                json={"contents": [{"parts": [
                    {"inline_data": {"mime_type": "image/jpeg", "data": b64}},
                    {"text": question}
                ]}]}
            )
            data = r.json()
            if "candidates" in data:
                result = data["candidates"][0]["content"]["parts"][0]["text"]
                logger.info("✅ Image analyzed via Gemini")
                return result
            logger.error(f"Gemini vision: {data.get('error', data)}")
    except Exception as e:
        logger.error(f"Gemini vision error: {e}")

    return None


async def search_web(query: str) -> str:
    """Поиск в интернете через Serper API"""
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            response = await client.post(
                "https://google.serper.dev/search",
                headers={
                    "X-API-KEY": SERPER_KEY,
                    "Content-Type": "application/json"
                },
                json={
                    "q": query,
                    "gl": "ru",
                    "hl": "ru",
                    "num": 5
                }
            )
            
            data = response.json()
            
            result_parts = []
            
            # Ответ от Google (быстрые ответы)
            if "answerBox" in data:
                answer_box = data["answerBox"]
                if "answer" in answer_box:
                    result_parts.append(f"📌 {answer_box['answer']}")
                elif "snippet" in answer_box:
                    result_parts.append(f"📌 {answer_box['snippet']}")
            
            # Органические результаты
            if "organic" in data:
                for item in data["organic"][:4]:
                    title = item.get("title", "")
                    snippet = item.get("snippet", "")
                    if title and snippet:
                        result_parts.append(f"• *{title}*\n{snippet}")
                    elif snippet:
                        result_parts.append(f"• {snippet}")
            
            if result_parts:
                return "\n\n".join(result_parts)
            else:
                return "По вашему запросу ничего не найдено."
                
    except Exception as e:
        logger.error(f"Search error: {e}")
        return f"Ошибка при поиске: {e}"

# =============================================
#  REMINDERS
# =============================================
def parse_reminder_time(text: str) -> Tuple[Optional[datetime], str]:
    """Парсит время из текста напоминания"""
    now = datetime.now()
    text_lower = text.lower()
    
    # Шаблоны для поиска времени
    patterns = [
        # через X минут
        (r'через (\d+)\s*минут', lambda m: now + timedelta(minutes=int(m.group(1)))),
        # через X час/часа/часов
        (r'через (\d+)\s*час', lambda m: now + timedelta(hours=int(m.group(1)))),
        # через X день/дня/дней
        (r'через (\d+)\s*дн', lambda m: now + timedelta(days=int(m.group(1)))),
        # завтра в HH:MM
        (r'завтра в (\d{1,2}):(\d{2})', lambda m: (now + timedelta(days=1)).replace(
            hour=int(m.group(1)), minute=int(m.group(2)), second=0, microsecond=0)),
        # в HH:MM (сегодня или завтра)
        (r'в (\d{1,2}):(\d{2})', lambda m: parse_hhmm(now, int(m.group(1)), int(m.group(2)))),
        # через X часов Y минут
        (r'через (\d+)\s*час(?:ов)?\s*(\d+)\s*минут', 
         lambda m: now + timedelta(hours=int(m.group(1)), minutes=int(m.group(2)))),
    ]
    
    for pattern, func in patterns:
        match = re.search(pattern, text_lower)
        if match:
            try:
                reminder_time = func(match)
                # Убираем временную часть из текста
                clean_text = re.sub(pattern, '', text, flags=re.IGNORECASE).strip()
                return reminder_time, clean_text
            except Exception as e:
                logger.error(f"Time parse error: {e}")
                continue
    
    return None, text

def parse_hhmm(now: datetime, hour: int, minute: int) -> datetime:
    """Парсит время в формате HH:MM"""
    candidate = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
    if candidate > now:
        return candidate
    else:
        return candidate + timedelta(days=1)

async def reminder_check_loop(app: Application) -> None:
    """Фоновый цикл проверки напоминаний"""
    while True:
        try:
            reminders = load_reminders()
            now = datetime.now()
            to_keep = []
            
            for reminder in reminders:
                reminder_time = datetime.fromisoformat(reminder["time"])
                if now >= reminder_time:
                    # Отправляем напоминание
                    try:
                        await app.bot.send_message(
                            chat_id=reminder["chat_id"],
                            text=f"⏰ *Напоминание!*\n\n{reminder['text']}",
                            parse_mode="Markdown"
                        )
                    except Exception as e:
                        logger.error(f"Failed to send reminder: {e}")
                else:
                    to_keep.append(reminder)
            
            if len(to_keep) != len(reminders):
                save_reminders(to_keep)
                
        except Exception as e:
            logger.error(f"Reminder loop error: {e}")
        
        await asyncio.sleep(60)  # Проверяем каждую минуту

# =============================================
#  MESSAGE SENDING
# =============================================
async def send_long_message(
    bot, 
    chat_id: int, 
    text: str, 
    parse_mode: str = "Markdown"
) -> None:
    """Отправляет длинное сообщение по частям"""
    max_length = 4096
    
    if len(text) <= max_length:
        try:
            await bot.send_message(
                chat_id=chat_id,
                text=text,
                parse_mode=parse_mode
            )
        except Exception:
            # Если Markdown не работает, отправляем без форматирования
            await bot.send_message(chat_id=chat_id, text=text)
        return
    
    # Разбиваем на части
    chunks = []
    for i in range(0, len(text), max_length - 500):
        chunk = text[i:i + max_length - 500]
        chunks.append(chunk)
    
    for i, chunk in enumerate(chunks, 1):
        header = f"📄 *Часть {i}/{len(chunks)}:*\n\n"
        full_text = header + chunk
        
        try:
            await bot.send_message(
                chat_id=chat_id,
                text=full_text,
                parse_mode=parse_mode
            )
        except Exception:
            await bot.send_message(chat_id=chat_id, text=full_text)
        
        await asyncio.sleep(0.5)  # Небольшая задержка между частями

# =============================================
#  COMMAND HANDLERS
# =============================================
async def start_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /start — всегда отвечает"""
    try:
        data = load_data()
        if "topics" not in data:
            data["topics"] = []
        chat = data["chats"].get(data.get("current", "main"), {"name": "Основной", "history": []})
        
        profile = data.get("profile", {})
        user_name = profile.get("name", "")
        greet_name = user_name if user_name else "господин"
        
        now_local = datetime.utcnow() + timedelta(hours=data.get("timezone_offset", 5))
        hour = now_local.hour
        if hour < 6: greeting_time = "🌙 Доброй ночи"
        elif hour < 12: greeting_time = "🌅 Доброе утро"
        elif hour < 17: greeting_time = "☀️ Добрый день"
        else: greeting_time = "🌆 Добрый вечер"
        
        active_tasks = sum(1 for t in data.get("tasks", []) if not t.get("done"))
        active_goals = sum(1 for g in data.get("goals", []) if g.get("status") == "active")
        habits_count = len(data.get("habits", []))
        
        text = (
            f"👑 *{greeting_time}, {greet_name}!*\n\n"
            f"💬 Чат: *{chat['name']}* | 📁 Всего: {len(data['chats'])}\n\n"
            f"📊 *Твой статус:*\n"
            f"├ ✅ Задач активных: *{active_tasks}*\n"
            f"├ 🎯 Целей в работе: *{active_goals}*\n"
            f"├ 🏃 Привычек: *{habits_count}*\n"
            f"└ 📚 Тем изучено: *{len(data['topics'])}*\n\n"
            f"🎙 *Нажми кнопку ниже для голосового агента*\n"
            f"_или просто напиши сообщение_"
        )
        await update.message.reply_text(text, parse_mode="Markdown", reply_markup=main_keyboard())
    except Exception as e:
        logger.error(f"start_command error: {e}")
        try:
            await update.message.reply_text(
                "👑 Приветствую, господин! Готов служить 👇",
                reply_markup=main_keyboard()
            )
        except Exception as e2:
            logger.error(f"start_command fallback error: {e2}")


async def restart_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /restart"""
    try:
        context.user_data.clear()
        await update.message.reply_text(
            "🔄 *Перезапуск выполнен, господин!*\n\nВсё готово к работе 👇",
            parse_mode="Markdown",
            reply_markup=main_keyboard()
        )
    except Exception as e:
        logger.error(f"restart error: {e}")
        await update.message.reply_text("🔄 Перезапущен!", reply_markup=main_keyboard())


async def memory_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /memory"""
    data = load_data()
    current_chat = data["chats"][data["current"]]
    
    # Последние темы
    topics_text = "\n".join([f"• {t}" for t in data["topics"][-15:]]) if data["topics"] else "Пока нет изученных тем"
    
    # Последние сообщения
    recent_messages = []
    for msg in current_chat["history"][-10:]:
        role_icon = "👤" if msg["role"] == "Пользователь" else "🤖"
        recent_messages.append(f"{role_icon} *{msg['time']}*\n{msg['text'][:100]}...")
    
    messages_text = "\n\n".join(recent_messages) if recent_messages else "Нет сообщений"
    
    memory_text = (
        f"🧠 *Память бота*\n\n"
        f"*Текущий чат:* {current_chat['name']}\n"
        f"*Всего сообщений:* {len(current_chat['history'])}\n"
        f"*Всего чатов:* {len(data['chats'])}\n\n"
        f"📚 *Изученные темы (последние 15):*\n{topics_text}\n\n"
        f"💬 *Последние сообщения:*\n\n{messages_text}"
    )
    
    await send_long_message(context.bot, update.effective_chat.id, memory_text)

async def clear_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /clear"""
    data = load_data()
    chat_name = data["chats"][data["current"]]["name"]
    
    # Очищаем историю текущего чата
    data["chats"][data["current"]]["history"] = []
    data["last_code"] = ""
    save_data(data)
    
    await update.message.reply_text(
        f"🗑 Чат *{chat_name}* полностью очищен!",
        parse_mode="Markdown",
        reply_markup=main_keyboard()
    )

async def status_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /status"""
    status_msg = await update.message.reply_text(
        "📡 Проверяю все AI провайдеры... Это займет около 8 секунд."
    )
    
    results = await check_all_providers()
    
    status_text = "📡 *Статус AI провайдеров:*\n\n"
    for name, status in results.items():
        status_text += f"{name}: {status}\n"
    
    await status_msg.edit_text(
        status_text,
        parse_mode="Markdown",
        reply_markup=main_keyboard()
    )

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Красивая справка по боту"""
    data = load_data()
    name = data.get("profile", {}).get("name", "господин")
    text = (
        f"👑 *Привет, {name}! Вот всё что я умею:*\n\n"
        
        "🎙 *Голос и чат*\n"
        "├ Голосовые сообщения — просто отправь\n"
        "├ Фото — анализирую что на нём\n"
        "└ Текст — отвечаю на любой вопрос\n\n"
        
        "✅ *Продуктивность*\n"
        "├ /tasks — задачи (добавь, отметь, удали)\n"
        "├ /goals — цели с прогресс-баром\n"
        "├ /habits — трекер привычек\n"
        "└ /reminders — напоминания в любое время\n\n"
        
        "🧠 *Знания и обучение*\n"
        "├ /brain — база знаний\n"
        "├ /tutor — репетитор по любой теме\n"
        "├ /research — глубокое исследование\n"
        "└ /books — книжный клуб\n\n"
        
        "🛠 *Инструменты*\n"
        "├ /weather — погода Ташкент\n"
        "├ /summarize [ссылка] — пересказ сайта/YouTube\n"
        "├ /search [запрос] — поиск в интернете\n"
        "├ /agent — автономный агент для сложных задач\n"
        "└ /export — экспорт в PDF\n\n"
        
        "⚙️ *Настройки*\n"
        "├ /profile — имя, город, о себе\n"
        "├ /timezone — часовой пояс (сейчас рассылка в 7:00)\n"
        "└ /style — персонализация стиля ответов\n\n"
        
        "💡 *Просто напиши* что хочешь — я пойму!"
    )
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("✅ Задачи", callback_data="tasks_show"),
         InlineKeyboardButton("🎯 Цели", callback_data="goals_show")],
        [InlineKeyboardButton("🌤 Погода", callback_data="weather_now"),
         InlineKeyboardButton("⏰ Напоминание", callback_data="reminder_add")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=kb)


async def new_chat_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /new"""
    await update.message.reply_text(
        "✏️ Введите название для нового чата:"
    )
    context.user_data["waiting_for"] = "new_chat_name"

async def chats_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /chats"""
    data = load_data()
    current = data["current"]
    
    keyboard = []
    for chat_id, chat_info in data["chats"].items():
        mark = "✅ " if chat_id == current else "💬 "
        button_text = f"{mark}{chat_info['name']} ({len(chat_info['history'])} сообщ.)"
        keyboard.append([InlineKeyboardButton(button_text, callback_data=f"switch_{chat_id}")])
    
    keyboard.append([
        InlineKeyboardButton("➕ Новый чат", callback_data="new_chat"),
        InlineKeyboardButton("🗑 Удалить", callback_data="delete_menu")
    ])
    
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await update.message.reply_text(
        "💬 *Ваши чаты:*\n\nВыберите чат для переключения:",
        parse_mode="Markdown",
        reply_markup=reply_markup
    )

async def search_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /search"""
    await update.message.reply_text(
        "🌐 Введите запрос для поиска в интернете:"
    )
    context.user_data["waiting_for"] = "web_search"

async def reminders_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /reminders"""
    chat_id = update.effective_chat.id
    all_reminders = load_reminders()
    user_reminders = [r for r in all_reminders if r["chat_id"] == chat_id]
    
    keyboard = []
    
    if user_reminders:
        text = "⏰ *Ваши напоминания:*\n\n"
        for i, rem in enumerate(user_reminders[:10]):  # Показываем последние 10
            rem_time = datetime.fromisoformat(rem["time"])
            time_str = rem_time.strftime("%d.%m.%Y в %H:%M")
            text += f"• {rem['text']} — *{time_str}*\n"
            keyboard.append([InlineKeyboardButton(
                f"🗑 Удалить: {rem['text'][:30]}...",
                callback_data=f"delete_reminder_{i}"
            )])
    else:
        text = (
            "⏰ *Напоминания*\n\n"
            "У вас пока нет активных напоминаний.\n\n"
            "*Примеры:*\n"
            "• _напомни через 30 минут позвонить_\n"
            "• _напомни в 15:30 встреча_\n"
            "• _напомни завтра в 10:00 купить хлеб_"
        )
    
    keyboard.append([InlineKeyboardButton("➕ Добавить напоминание", callback_data="add_reminder")])
    
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await update.message.reply_text(
        text,
        parse_mode="Markdown",
        reply_markup=reply_markup
    )

async def templates_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик кнопки Шаблоны"""
    template_items = list(TEMPLATES.items())
    
    keyboard = []
    # Создаем ряды по 2 кнопки
    for i in range(0, len(template_items), 2):
        row = []
        # Первая кнопка в ряду
        row.append(InlineKeyboardButton(
            template_items[i][0],
            callback_data=f"template_{i}"
        ))
        # Вторая кнопка, если есть
        if i + 1 < len(template_items):
            row.append(InlineKeyboardButton(
                template_items[i + 1][0],
                callback_data=f"template_{i + 1}"
            ))
        keyboard.append(row)
    
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await update.message.reply_text(
        "📋 *Выберите шаблон запроса:*\n\n"
        "После выбора скопируйте текст и дополните его.",
        parse_mode="Markdown",
        reply_markup=reply_markup
    )

# =============================================
#  CALLBACK QUERY HANDLER
# =============================================
async def button_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик нажатий на инлайн кнопки"""
    query = update.callback_query
    await query.answer()
    
    data = load_data()
    chat_id = update.effective_chat.id
    
    # Переключение чата
    if data.startswith("switch_"):
        chat_key = data[7:]
        if chat_key in data["chats"]:
            data["current"] = chat_key
            save_data(data)
            await query.edit_message_text(
                f"✅ Переключился на чат: *{data['chats'][chat_key]['name']}*",
                parse_mode="Markdown"
            )
    
    # Меню удаления чатов
    elif data == "delete_menu":
        keyboard = []
        for chat_id_key, chat_info in data["chats"].items():
            if chat_id_key != "main":  # Не даем удалить основной чат
                keyboard.append([InlineKeyboardButton(
                    f"🗑 {chat_info['name']}",
                    callback_data=f"delete_{chat_id_key}"
                )])
        
        if keyboard:
            keyboard.append([InlineKeyboardButton("◀️ Назад", callback_data="back_to_chats")])
            reply_markup = InlineKeyboardMarkup(keyboard)
            await query.edit_message_text(
                "🗑 *Выберите чат для удаления:*\n\nОсновной чат удалить нельзя.",
                parse_mode="Markdown",
                reply_markup=reply_markup
            )
        else:
            await query.edit_message_text(
                "❌ Нет чатов для удаления, кроме основного.",
                reply_markup=InlineKeyboardMarkup([[
                    InlineKeyboardButton("◀️ Назад", callback_data="back_to_chats")
                ]])
            )
    
    # Удаление конкретного чата
    elif data.startswith("delete_"):
        chat_key = data[7:]
        if chat_key in data["chats"] and chat_key != "main":
            chat_name = data["chats"][chat_key]["name"]
            del data["chats"][chat_key]
            
            # Если удалили текущий чат, переключаемся на основной
            if data["current"] == chat_key:
                data["current"] = "main"
            
            save_data(data)
            await query.edit_message_text(
                f"🗑 Чат *{chat_name}* успешно удален!",
                parse_mode="Markdown"
            )
        else:
            await query.edit_message_text(
                "❌ Нельзя удалить основной чат!"
            )
    
    # Назад к списку чатов
    elif data == "back_to_chats":
        current = data["current"]
        keyboard = []
        for chat_id_key, chat_info in data["chats"].items():
            mark = "✅ " if chat_id_key == current else "💬 "
            button_text = f"{mark}{chat_info['name']} ({len(chat_info['history'])} сообщ.)"
            keyboard.append([InlineKeyboardButton(button_text, callback_data=f"switch_{chat_id_key}")])
        
        keyboard.append([
            InlineKeyboardButton("➕ Новый чат", callback_data="new_chat"),
            InlineKeyboardButton("🗑 Удалить", callback_data="delete_menu")
        ])
        
        reply_markup = InlineKeyboardMarkup(keyboard)
        await query.edit_message_text(
            "💬 *Ваши чаты:*\n\nВыберите чат для переключения:",
            parse_mode="Markdown",
            reply_markup=reply_markup
        )
    
    # Создание нового чата
    elif data == "new_chat":
        await query.edit_message_text(
            "✏️ Введите название для нового чата:"
        )
        context.user_data["waiting_for"] = "new_chat_name"
    
    # Добавление напоминания
    elif data == "add_reminder":
        await query.edit_message_text(
            "⏰ *Создание напоминания*\n\n"
            "Напишите текст напоминания с указанием времени.\n\n"
            "*Примеры:*\n"
            "• _напомни через 30 минут позвонить маме_\n"
            "• _напомни в 15:30 встреча с клиентом_\n"
            "• _напомни завтра в 10:00 купить продукты_",
            parse_mode="Markdown"
        )
        context.user_data["waiting_for"] = "reminder"
    
    # Удаление напоминания
    elif data.startswith("delete_reminder_"):
        try:
            idx = int(data[15:])
            all_reminders = load_reminders()
            user_reminders = [r for r in all_reminders if r["chat_id"] == chat_id]
            
            if idx < len(user_reminders):
                reminder_to_delete = user_reminders[idx]
                all_reminders.remove(reminder_to_delete)
                save_reminders(all_reminders)
                
                await query.edit_message_text(
                    f"🗑 Напоминание *{reminder_to_delete['text']}* удалено!",
                    parse_mode="Markdown"
                )
            else:
                await query.edit_message_text("❌ Напоминание не найдено.")
        except Exception as e:
            logger.error(f"Delete reminder error: {e}")
            await query.edit_message_text("❌ Ошибка при удалении напоминания.")
    
    # Шаблон запроса
    elif data.startswith("template_"):
        try:
            idx = int(data[9:])
            template_keys = list(TEMPLATES.keys())
            if idx < len(template_keys):
                key = template_keys[idx]
                template_text = TEMPLATES[key]
                await query.edit_message_text(
                    f"📋 *{key}*\n\n"
                    f"Скопируйте и дополните запрос:\n\n"
                    f"`{template_text}`",
                    parse_mode="Markdown"
                )
        except Exception as e:
            logger.error(f"Template error: {e}")
            await query.edit_message_text("❌ Ошибка при загрузке шаблона.")

# =============================================
#  MESSAGE HANDLERS
# =============================================
async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик текстовых сообщений"""
    text = update.message.text
    chat_id = update.effective_chat.id
    
    # Обработка кнопок главного меню
    if text == "💬 Новый чат":
        await new_chat_command(update, context)
        return
    elif text == "📂 Мои чаты":
        await chats_command(update, context)
        return
    elif text == "📋 Шаблоны":
        await templates_command(update, context)
        return
    elif text == "🔍 Поиск":
        await search_command(update, context)
        return
    elif text == "⏰ Напоминания":
        await reminders_command(update, context)
        return
    elif text == "✅ Задачи":
        await tasks_command(update, context)
        return
    elif text == "👤 Профиль":
        await profile_command(update, context)
        return
    elif text == "📊 Статистика":
        await stats_command(update, context)
        return
    elif text == "🎯 Цели":
        await goals_command(update, context)
        return
    elif text == "📚 База знаний":
        await brain_command(update, context)
        return
    elif text == "🔬 Исследовать":
        await update.message.reply_text("🔬 Напиши тему для глубокого исследования:")
        context.user_data["waiting_for"] = "deep_research"
        return
    elif text == "🎓 Репетитор":
        await tutor_command(update, context)
        return
    elif text == "🤝 Дебаты":
        await debate_command(update, context)
        return
    elif text == "💼 Собеседование":
        await interview_command(update, context)
        return
    elif text == "📖 Книги":
        await books_command(update, context)
        return
    elif text == "🌐 Мониторинг":
        await monitor_command(update, context)
        return
    elif text == "📧 Email":
        await email_command(update, context)
        return
    elif text in ("📍 Рядом", "📍 Найти рядом"):
        keyboard_loc = InlineKeyboardMarkup([
            [InlineKeyboardButton("📍 Отправить геолокацию", callback_data="request_location")]
        ])
        await update.message.reply_text(
            "📍 *Поиск мест рядом*\n\nНажми кнопку ниже — отправь геолокацию\nИли напиши: _найди кафе_, _где McDonald's_, _музеи рядом_",
            parse_mode="Markdown",
            reply_markup=ReplyKeyboardMarkup(
                [[KeyboardButton("📍 Отправить мою геолокацию", request_location=True)]],
                one_time_keyboard=True, resize_keyboard=True
            )
        )
        return
    elif text == "🧠 Мой стиль":
        await finetune_command(update, context)
        return
    elif text == "🤖 Агент":
        await agent_command(update, context)
        return
    elif text == "📤 Экспорт PDF":
        await export_command(update, context)
        return
    elif text == "🏃 Привычки":
        await habits_command(update, context)
        return
    elif text == "📈 График из файла":
        await update.message.reply_text(
            "📈 *График из файла*\n\n"
            "Отправь Excel (.xlsx) или CSV файл с подписью _график_\n\n"
            "Я создам красивый график из твоих данных!",
            parse_mode="Markdown"
        )
        return
    elif text == "🧠 Память":
        await memory_command(update, context)
        return
    elif text == "🌤 Погода":
        await weather_command(update, context)
        return
    elif text in ("🔗 Ссылка", "🔗 Пересказ ссылки"):
        await summarize_command(update, context)
        return
    elif text in ("📄 Файл", "📄 Анализ файла"):
        await update.message.reply_text(
            "Отправь файл (.txt, .pdf, .docx, .csv)\n\n"
            "Добавь подпись к файлу:\n"
            "конспект — краткое содержание\n"
            "викторина — тест с вопросами\n"
            "или без подписи для общего анализа"
        )
        return
    elif text == "❓ Помощь":
        await help_command(update, context)
        return
    
    # Проверяем, ждем ли мы какой-то ввод
    waiting_for = context.user_data.get("waiting_for")
    
    if waiting_for == "new_chat_name":
        # Создаем новый чат
        data = load_data()
        chat_key = f"chat_{int(datetime.now().timestamp())}"
        data["chats"][chat_key] = {
            "name": text,
            "history": [],
            "created": datetime.now().strftime("%d.%m.%Y %H:%M")
        }
        data["current"] = chat_key
        save_data(data)
        
        context.user_data.pop("waiting_for", None)
        
        await update.message.reply_text(
            f"✅ Создан новый чат: *{text}*",
            parse_mode="Markdown",
            reply_markup=main_keyboard()
        )
        return
    
    elif waiting_for == "weather_city":
        context.user_data.pop("waiting_for", None)
        msg = await update.message.reply_text(f"🌤 Ищу погоду для *{text}*...", parse_mode="Markdown")
        data = load_data()
        data["weather_city"] = text
        save_data(data)
        weather = await get_weather(text)
        await msg.edit_text(weather, parse_mode="Markdown")
        return

    elif waiting_for == "summarize_url":
        context.user_data.pop("waiting_for", None)
        import re
        url_match = re.search(r'https?://\S+', text)
        if not url_match:
            await update.message.reply_text("❌ Не нашёл ссылку. Отправь корректный URL начинающийся с http://")
            return
        url = url_match.group(0)
        msg = await update.message.reply_text(f"🔗 Читаю материал...", parse_mode="Markdown")
        result = await summarize_url(url)
        await msg.delete()
        await send_long_message(context.application.bot, chat_id, result, parse_mode="Markdown")
        return

    elif waiting_for == "generate_prompt":
        context.user_data.pop("waiting_for", None)
        msg = await update.message.reply_text("✍️ Создаю промпты...")
        result = await generate_prompt(text)
        await msg.delete()
        await send_long_message(context.application.bot, chat_id, result, parse_mode="Markdown")
        return

    elif waiting_for == "profile_name":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("profile", {})["name"] = text
        save_data(d)
        await update.message.reply_text(f"✅ Имя обновлено: *{text}*", parse_mode="Markdown")
        return
    elif waiting_for == "profile_city":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("profile", {})["city"] = text
        d["weather_city"] = text
        save_data(d)
        await update.message.reply_text(f"✅ Город обновлён: *{text}*", parse_mode="Markdown")
        return
    elif waiting_for == "profile_interest":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("profile", {}).setdefault("interests", []).append(text)
        save_data(d)
        await update.message.reply_text(f"✅ Интерес добавлен: *{text}*", parse_mode="Markdown")
        return
    elif waiting_for == "profile_goal":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("profile", {}).setdefault("goals", []).append(text)
        save_data(d)
        await update.message.reply_text(f"✅ Цель добавлена: *{text}*", parse_mode="Markdown")
        return
    elif waiting_for == "deep_research":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        asyncio.create_task(deep_research(context.application, chat_id, text, d))
        return
    elif waiting_for == "tutor_topic":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        asyncio.create_task(create_tutor_course(context.application, chat_id, text, d))
        return
    elif waiting_for == "goal_add":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("goals_system", {"goals": []})["goals"].append({
            "title": text, "status": "active", "progress": 0,
            "created": datetime.now().strftime("%d.%m.%Y"),
            "deadline": "", "steps": []
        })
        save_data(d)
        await update.message.reply_text(f"🎯 Цель добавлена: *{text}*\n\nТеперь напиши дедлайн (например: 01.06.2025) или 'без срока':", parse_mode="Markdown")
        context.user_data["waiting_for"] = "goal_deadline"
        context.user_data["last_goal"] = text
        return
    elif waiting_for == "nearby_custom_query":
        context.user_data.pop("waiting_for", None)
        loc = USER_LOCATIONS.get(chat_id)
        if not loc:
            await update.message.reply_text("❌ Геолокация не найдена. Отправь снова через 📍 Найти рядом")
            return
        await update.message.reply_text(f"🔍 Ищу: *{text}*...", parse_mode="Markdown")
        result = await search_nearby_nominatim(loc["lat"], loc["lon"], text)
        await update.message.reply_text(result, parse_mode="Markdown", disable_web_page_preview=True)
        return

    elif waiting_for == "email_address":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("email_config", {})["address"] = text
        save_data(d)
        await update.message.reply_text(
            f"✅ Email сохранён: *{text}*\n\nТеперь введи App Password от Google:",
            parse_mode="Markdown"
        )
        context.user_data["waiting_for"] = "email_password"
        return

    elif waiting_for == "email_password":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("email_config", {})["password"] = text
        save_data(d)
        await update.message.reply_text(
            "✅ Email подключён! Проверяю подключение...",
            parse_mode="Markdown"
        )
        cfg = d["email_config"]
        emails = await fetch_emails(cfg["address"], cfg["password"], 3)
        if emails and "error" not in emails[0]:
            await update.message.reply_text(
                f"🎉 Подключение успешно!\n\n📬 Непрочитанных: {len(emails)}\n\nНапиши /email чтобы читать почту.",
                parse_mode="Markdown"
            )
        else:
            err = emails[0].get("error", "неизвестная ошибка") if emails else "нет данных"
            await update.message.reply_text(f"❌ Ошибка подключения: {err}\n\nПроверь App Password и попробуй снова.")
        return

    elif waiting_for == "email_compose_to":
        context.user_data["email_to"] = text
        context.user_data["waiting_for"] = "email_compose_subject"
        await update.message.reply_text("📝 Тема письма:")
        return

    elif waiting_for == "email_compose_subject":
        context.user_data["email_subject"] = text
        context.user_data["waiting_for"] = "email_compose_body"
        await update.message.reply_text("✉️ Текст письма:")
        return

    elif waiting_for == "email_compose_body":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        cfg = d.get("email_config", {})
        if not cfg.get("address"):
            await update.message.reply_text("❌ Email не подключён. Используй /email")
            return
        to = context.user_data.pop("email_to", "")
        subject = context.user_data.pop("email_subject", "")
        # AI улучшает письмо
        improve_prompt = f"Улучши это письмо, сделай его профессиональным:\n\nКому: {to}\nТема: {subject}\nТекст: {text}"
        improved, _ = await get_best_ai_response(improve_prompt, "Ты помощник по деловой переписке.")
        body = improved or text
        ok = await send_email(cfg["address"], cfg["password"], to, subject, body)
        if ok:
            await update.message.reply_text(f"✅ Письмо отправлено!\n\n📬 Кому: {to}\n📋 Тема: {subject}", parse_mode="Markdown")
        else:
            await update.message.reply_text("❌ Ошибка отправки. Проверь настройки.")
        return

    elif waiting_for == "agent_task":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        asyncio.create_task(run_agent(context.application, chat_id, text, d))
        return

    elif waiting_for == "debate_topic":
        context.user_data.pop("waiting_for", None)
        context.user_data["debate_topic"] = text
        context.user_data["debate_round"] = 0
        DEBATE_SESSIONS[chat_id] = {"topic": text, "round": 0}
        await update.message.reply_text(
            f"🤝 *Дебаты: {text}*\n\nТы за, я против!\nНапиши свой первый аргумент:",
            parse_mode="Markdown"
        )
        context.user_data["waiting_for"] = "debate_arg"
        return

    elif waiting_for == "debate_arg":
        session = DEBATE_SESSIONS.get(chat_id, {})
        topic = session.get("topic", context.user_data.get("debate_topic", ""))
        round_num = session.get("round", 0) + 1
        DEBATE_SESSIONS[chat_id] = {"topic": topic, "round": round_num}
        context.user_data["waiting_for"] = "debate_arg"
        asyncio.create_task(run_debate(context.application, chat_id, topic, text, round_num))
        return

    elif waiting_for == "interview_role":
        context.user_data.pop("waiting_for", None)
        INTERVIEW_SESSIONS[chat_id] = {"role": text, "q_num": 1}
        context.user_data["waiting_for"] = "interview_answer"
        asyncio.create_task(run_interview(context.application, chat_id, text, "", 1))
        return

    elif waiting_for == "interview_answer":
        session = INTERVIEW_SESSIONS.get(chat_id, {})
        role = session.get("role", "")
        q_num = session.get("q_num", 1) + 1
        INTERVIEW_SESSIONS[chat_id] = {"role": role, "q_num": q_num}
        if q_num <= 8:
            context.user_data["waiting_for"] = "interview_answer"
        else:
            context.user_data.pop("waiting_for", None)
        asyncio.create_task(run_interview(context.application, chat_id, role, text, q_num))
        return

    elif waiting_for == "book_add":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("books", []).append({
            "title": text, "author": "", "status": "want",
            "added": datetime.now().strftime("%d.%m.%Y"), "rating": 0
        })
        save_data(d)
        asyncio.create_task(analyze_book(context.application, chat_id, text))
        return

    elif waiting_for == "monitor_add":
        context.user_data.pop("waiting_for", None)
        import re as _re3
        url_m = _re3.search(r'https?://\S+', text)
        url = url_m.group(0) if url_m else text
        monitors = load_monitors()
        monitors.append({
            "url": url, "chat_id": chat_id,
            "active": True, "content_hash": "",
            "added": datetime.now().strftime("%d.%m.%Y %H:%M"),
            "last_check": ""
        })
        save_monitors(monitors)
        await update.message.reply_text(
            f"✅ Мониторинг добавлен!\n\n🌐 {url}\n\nБуду проверять каждые 30 минут и сообщу об изменениях.",
            parse_mode="Markdown"
        )
        return

    elif waiting_for == "book_analyze":
        context.user_data.pop("waiting_for", None)
        asyncio.create_task(analyze_book(context.application, chat_id, text))
        return

    elif waiting_for == "habit_add":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("habits_tracker", []).append({
            "name": text.strip(),
            "created": datetime.now().strftime("%d.%m.%Y"),
            "checks": {}
        })
        save_data(d)
        await update.message.reply_text(
            f"✅ Привычка добавлена: *{text}*\n\nОтмечай каждый день через кнопку 🏃 Привычки!",
            parse_mode="Markdown"
        )
        return

    elif waiting_for == "goal_deadline":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        goals = d.get("goals_system", {}).get("goals", [])
        last = context.user_data.pop("last_goal", "")
        for g in goals:
            if g["title"] == last:
                g["deadline"] = text if text.lower() != "без срока" else "без срока"
        save_data(d)
        await update.message.reply_text("✅ Цель сохранена!\n\nНапиши /goals чтобы посмотреть и обновлять прогресс.", parse_mode="Markdown")
        return
    elif waiting_for == "goal_update_progress":
        context.user_data.pop("waiting_for", None)
        idx = context.user_data.pop("goal_update_idx", 0)
        try:
            progress = max(0, min(100, int(text.strip().replace("%",""))))
        except:
            progress = 0
        d = load_data()
        goals = d.get("goals_system", {}).get("goals", [])
        if 0 <= idx < len(goals):
            goals[idx]["progress"] = progress
            if progress >= 100:
                goals[idx]["status"] = "done"
                goals[idx]["done_at"] = datetime.now().strftime("%d.%m.%Y")
                await update.message.reply_text(f"🎉 *Цель выполнена: {goals[idx]['title']}*", parse_mode="Markdown")
            else:
                await update.message.reply_text(f"✅ Прогресс обновлён: *{progress}%*", parse_mode="Markdown")
        save_data(d)
        return

    elif waiting_for == "brain_search":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        knowledge = d.get("second_brain", {}).get("knowledge", {})
        query = text.lower()
        found = [(k, v) for k, v in knowledge.items() if query in k or query in str(v).lower()]
        if found:
            lines = ["🔍 *Найдено в базе знаний:*\n"]
            for k, v in found[:5]:
                lines.append(f"📖 *{v['title']}*")
                if v.get("notes"):
                    lines.append(v["notes"][-1]["text"][:200])
                lines.append("")
            await update.message.reply_text("\n".join(lines), parse_mode="Markdown")
        else:
            await update.message.reply_text(f"❌ По запросу '{text}' ничего не найдено в базе знаний.")
        return
    elif waiting_for == "task_add":
        context.user_data.pop("waiting_for", None)
        d = load_data()
        d.setdefault("tasks", []).append({
            "text": text,
            "done": False,
            "created": datetime.now().strftime("%d.%m.%Y %H:%M"),
            "priority": "medium"
        })
        save_data(d)
        await update.message.reply_text(f"✅ Задача добавлена: *{text}*", parse_mode="Markdown")
        return
    elif waiting_for == "web_search":
        # Выполняем поиск
        context.user_data.pop("waiting_for", None)
        
        search_msg = await update.message.reply_text("🌐 Ищу в интернете...")
        
        # Поиск в интернете
        search_results = await search_web(text)
        
        # Отправляем результаты
        await search_msg.edit_text(
            f"🌐 *Результаты поиска:*\n\n{search_results[:3000]}",
            parse_mode="Markdown"
        )
        
        # Дополнительно спрашиваем AI для краткого ответа
        data = load_data()
        system_context = build_context(data) + f"\n\nРезультаты поиска:\n{search_results}"
        
        thinking = await update.message.reply_text("🤔 Анализирую результаты...")
        
        answer, ai_name = await get_best_ai_response(
            f"На основе результатов поиска дай краткий ответ на вопрос: {text}",
            system_context
        )
        
        if answer:
            await thinking.edit_text(
                f"💡 *Ответ ({ai_name}):*\n\n{answer[:3000]}",
                parse_mode="Markdown"
            )
        else:
            await thinking.delete()
        
        return
    
    elif waiting_for == "reminder":
        # Создаем напоминание
        context.user_data.pop("waiting_for", None)
        
        reminder_time, reminder_text = parse_reminder_time(text)
        
        if reminder_time:
            reminders = load_reminders()
            reminders.append({
                "chat_id": chat_id,
                "text": reminder_text,
                "time": reminder_time.isoformat()
            })
            save_reminders(reminders)
            
            time_str = reminder_time.strftime("%d.%m.%Y в %H:%M")
            await update.message.reply_text(
                f"⏰ *Напоминание установлено!*\n\n"
                f"📝 {reminder_text}\n"
                f"🕐 {time_str}",
                parse_mode="Markdown",
                reply_markup=main_keyboard()
            )
        else:
            await update.message.reply_text(
                "❌ Не удалось распознать время.\n\n"
                "Используйте формат:\n"
                "• _через 30 минут ..._\n"
                "• _в 15:30 ..._\n"
                "• _завтра в 10:00 ..._",
                parse_mode="Markdown",
                reply_markup=main_keyboard()
            )
        return
    
    # Обычное сообщение
    # Проверяем, не напоминание ли это
    if any(word in text.lower() for word in ["напомни", "напоминание"]):
        reminder_time, reminder_text = parse_reminder_time(text)
        if reminder_time:
            reminders = load_reminders()
            reminders.append({
                "chat_id": chat_id,
                "text": reminder_text,
                "time": reminder_time.isoformat()
            })
            save_reminders(reminders)
            
            time_str = reminder_time.strftime("%d.%m.%Y в %H:%M")
            await update.message.reply_text(
                f"⏰ *Напоминание установлено!*\n\n"
                f"📝 {reminder_text}\n"
                f"🕐 {time_str}",
                parse_mode="Markdown"
            )
            return
    
    # Авто-определение поиска рядом по тексту
    _text_low3 = text.lower()
    _nearby_keywords = ["найди рядом", "где рядом", "поблизости", "рядом со мной",
                        "ближайший", "ближайшая", "ближайшие", "музей рядом",
                        "кафе рядом", "ресторан рядом", "аптека рядом"]
    _search_keywords = ["найди ", "где находится", "найти ", "покажи "]

    if any(k in _text_low3 for k in _nearby_keywords):
        loc = USER_LOCATIONS.get(chat_id)
        if loc:
            # Определяем категорию
            if any(w in _text_low3 for w in ["кафе", "ресторан", "поесть", "покушать", "еда"]):
                cat, label = "food", "🍽 Кафе и рестораны"
            elif any(w in _text_low3 for w in ["музей", "галерея", "достопримечательност"]):
                cat, label = "museum", "🏛 Музеи"
            elif any(w in _text_low3 for w in ["аптека", "лекарств"]):
                cat, label = "pharmacy", "💊 Аптеки"
            elif any(w in _text_low3 for w in ["банк", "банком", "atm", "банкомат"]):
                cat, label = "bank", "🏦 Банки"
            elif any(w in _text_low3 for w in ["магазин", "супермаркет", "продукт"]):
                cat, label = "shop", "🛒 Магазины"
            elif any(w in _text_low3 for w in ["отель", "гостиниц", "хостел"]):
                cat, label = "hotel", "🏨 Отели"
            elif any(w in _text_low3 for w in ["кофе", "кофейня"]):
                cat, label = "cafe", "☕ Кофейни"
            else:
                cat, label = "food", "📍 Места"
            result = await search_nearby(loc["lat"], loc["lon"], cat, label)
            await update.message.reply_text(result, parse_mode="Markdown", disable_web_page_preview=True)
        else:
            await update.message.reply_text(
                "📍 Сначала отправь геолокацию!\n\nНажми кнопку *📍 Найти рядом* в меню.",
                parse_mode="Markdown",
                reply_markup=ReplyKeyboardMarkup(
                    [[KeyboardButton("📍 Отправить мою геолокацию", request_location=True)]],
                    one_time_keyboard=True, resize_keyboard=True
                )
            )
        return

    # Поиск по названию заведения
    if any(k in _text_low3 for k in _search_keywords) and not any(k in _text_low3 for k in ["найди задачу", "найди файл", "найди чат"]):
        import re as _re_loc
        # Извлекаем название
        m_loc = _re_loc.search(r'(?:найди|найти|где находится|покажи)\s+(.+?)(?:\s+рядом|\s+поблизости|$)', _text_low3)
        if m_loc:
            search_name = m_loc.group(1).strip()
            loc = USER_LOCATIONS.get(chat_id)
            if loc and len(search_name) > 2:
                result = await search_nearby_nominatim(loc["lat"], loc["lon"], search_name)
                await update.message.reply_text(result, parse_mode="Markdown", disable_web_page_preview=True)
                return

    # Авто-определение маршрутов
    _is_route, _from, _to = is_route_request(text)
    if _is_route:
        route_info = await get_route_info(_from, _to)
        await update.message.reply_text(route_info, parse_mode="Markdown", disable_web_page_preview=True)
        return

    # Авто-определение агентного режима
    _text_low2 = text.lower()
    if any(p in _text_low2 for p in ["агент выполни", "агент сделай", "агент исследуй", "/agent "]):
        _task = text.split(" ", 1)[1] if " " in text else text
        _d = load_data()
        asyncio.create_task(run_agent(context.application, chat_id, _task, _d))
        return

    # Авто-определение исследования
    _is_research, _research_topic = is_research_request(text)
    if _is_research:
        _d = load_data()
        asyncio.create_task(deep_research(context.application, chat_id, _research_topic, _d))
        return

    # Авто-определение автопилота
    _is_auto, _auto_task = is_autopilot_request(text)
    if _is_auto:
        _d = load_data()
        asyncio.create_task(autopilot_task(context.application, chat_id, _auto_task, _d))
        return

    # Авто-определение задач в тексте
    _text_low = text.lower()
    if any(p in _text_low for p in ["добавь задачу", "задача ", "новая задача", "добавить задачу"]):
        import re as _re2
        task_match = _re2.search(r'(?:добавь задачу|задача|новая задача|добавить задачу)[:\s]+(.+)', _text_low)
        if task_match:
            task_text = task_match.group(1).strip().capitalize()
            data = load_data()
            data.setdefault("tasks", [])
            data["tasks"].append({
                "text": task_text,
                "done": False,
                "created": datetime.now().strftime("%d.%m.%Y %H:%M"),
                "priority": "medium"
            })
            save_data(data)
            active = len([t for t in data["tasks"] if not t.get("done")])
            await update.message.reply_text(
                f"✅ Задача добавлена!\n\n📋 *{task_text}*\n\nВсего активных задач: {active}",
                parse_mode="Markdown"
            )
            return

    # Авто-определение погоды
    import re as _re
    _wm = _re.search(r'погод[ауеыи]?\s+(?:в\s+)?([а-яёА-ЯЁ][а-яёА-ЯЁ\s]{1,25})', text.lower())
    if _wm or any(w in text.lower() for w in ["какая погода", "погода сегодня", "погода завтра", "прогноз погоды", "погода в"]):
        _city = _wm.group(1).strip().title() if _wm else load_data().get("weather_city", "Tashkent")
        _msg = await update.message.reply_text(f"🌤 Получаю погоду для {_city}...")
        _weather = await get_weather(_city)
        _d = load_data(); _d["weather_city"] = _city; save_data(_d)
        await _msg.edit_text(_weather, parse_mode="Markdown")
        return

    # Авто-определение ссылок
    _um = _re.search(r'https?://\S+', text)
    if _um:
        _url = _um.group(0)
        _msg = await update.message.reply_text("🔗 Читаю материал по ссылке...")
        _result = await summarize_url(_url)
        await _msg.delete()
        await send_long_message(context.application.bot, chat_id, _result, parse_mode="Markdown")
        return

    # Авто-определение запроса промпта
    # Авто-команда запомни
    if text.lower().startswith("запомни:") or text.lower().startswith("запомни "):
        fact = text[8:].strip() if text.lower().startswith("запомни:") else text[8:].strip()
        if fact:
            _d = load_data()
            _d.setdefault("second_brain", {}).setdefault("facts", []).append(fact)
            save_data(_d)
            await update.message.reply_text(f"💡 Запомнил: *{fact}*", parse_mode="Markdown")
            return

    if any(p in text.lower() for p in ["напиши промпт", "создай промпт", "промпт для", "prompt для"]):
        _msg = await update.message.reply_text("✍️ Создаю промпты...")
        _result = await generate_prompt(text)
        await _msg.delete()
        await send_long_message(context.application.bot, chat_id, _result, parse_mode="Markdown")
        return

    # Основная обработка через AI
    data = load_data()
    
    # Сообщаем, что начали работу
    await update.message.reply_text("⚙️ Работаю... Это может занять до 30 секунд.")
    
    # Запускаем обработку в фоне
    asyncio.create_task(process_message(
        context.application,
        chat_id,
        text,
        data,
        update.message.message_id
    ))

async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик голосовых сообщений"""
    chat_id = update.effective_chat.id
    
    voice_msg = await update.message.reply_text("🎤 Распознаю голосовое сообщение...")
    
    try:
        # Получаем файл голосового сообщения
        file = await context.bot.get_file(update.message.voice.file_id)
        
        # Скачиваем файл правильно
        audio_data = bytes(await file.download_as_bytearray())
        
        # Распознаем текст
        transcribed_text = await transcribe_voice(audio_data)
        
        if transcribed_text:
            await voice_msg.edit_text(
                f"🎤 *Распознано:*\n\n{transcribed_text}",
                parse_mode="Markdown"
            )
            
            # Обрабатываем распознанный текст
            data = load_data()
            asyncio.create_task(process_message(
                context.application,
                chat_id,
                transcribed_text,
                data,
                update.message.message_id
            ))
        else:
            await voice_msg.edit_text("❌ Не удалось распознать голосовое сообщение.")
            
    except Exception as e:
        logger.error(f"Voice processing error: {e}")
        await voice_msg.edit_text(f"❌ Ошибка при обработке голоса: {e}")

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик фотографий"""
    chat_id = update.effective_chat.id
    caption = update.message.caption or "Опиши подробно что видишь на этом изображении."
    
    photo_msg = await update.message.reply_text("📸 Анализирую изображение...")
    
    try:
        # Берем фото самого высокого качества
        photo = update.message.photo[-1]
        file = await context.bot.get_file(photo.file_id)
        
        # Скачиваем фото правильно
        image_data = bytes(await file.download_as_bytearray())
        
        # Анализируем через Gemini Vision
        analysis = await analyze_image(image_data, caption)
        
        if analysis:
            await photo_msg.edit_text(
                f"📸 *Анализ изображения:*\n\n{analysis}",
                parse_mode="Markdown"
            )
        else:
            await photo_msg.edit_text(
                "❌ Не удалось проанализировать изображение.\n\n"
                "Попробуй:\n"
                "• Отправить фото с подписью что именно нужно найти\n"
                "• Убедись что фото чёткое и не слишком тёмное"
            )
            
    except Exception as e:
        logger.error(f"Photo processing error: {e}")
        await photo_msg.edit_text(f"❌ Ошибка при анализе фото: {e}")

# =============================================
#  MESSAGE PROCESSING
# =============================================
async def process_message(
    app: Application,
    chat_id: int,
    text: str,
    data: Dict[str, Any],
    reply_to_message_id: int = None
) -> None:
    """Обрабатывает сообщение и генерирует ответ через AI"""
    
    try:
        # Обновляем статистику
        update_stats(data, "message")
        # Сохраняем сообщение пользователя
        add_message(data, "Пользователь", text)
        
        # Проверяем, нужно ли изучать тему
        if any(word in text.lower() for word in ["изучи", "расскажи", "объясни", "что такое", "кто такой"]):
            topic = text[:60]
            if topic not in data["topics"]:
                data["topics"].append(topic)
                asyncio.create_task(save_to_brain(data, topic, topic))
        
        # Строим контекст
        system_context = build_context(data)
        
        # Определяем тип запроса
        is_code = is_code_request(text)
        is_search = is_search_request(text)
        
        result = None
        ai_name = None
        
        if is_search:
            # Поиск в интернете
            search_results = await search_web(text)
            if search_results and search_results != "По вашему запросу ничего не найдено.":
                await app.bot.send_message(
                    chat_id=chat_id,
                    text=f"🌐 *Результаты поиска:*\n\n{search_results[:3000]}",
                    parse_mode="Markdown"
                )
            
            # Получаем ответ AI с учетом результатов поиска
            search_context = system_context + f"\n\nРезультаты поиска:\n{search_results}"
            result, ai_name = await get_best_ai_response(
                f"На основе результатов поиска дай ответ: {text}",
                search_context
            )
        
        elif is_code:
            # Генерация кода
            result, ai_name = await generate_code_smart(text, system_context)
        
        else:
            # Обычный текстовый ответ
            result, ai_name = await get_best_ai_response(text, system_context)
        
        if not result:
            await app.bot.send_message(
                chat_id=chat_id,
                text="❌ Ни один AI провайдер не ответил. Попробуйте позже."
            )
            return
        
        # Сохраняем ответ
        add_message(data, "Агент", result[:400])
        
        # Проверяем, есть ли код в ответе
        code = extract_code_from_text(result)
        
        if code and len(code) > 100:
            # Сохраняем код в память
            data["last_code"] = code[:800]
            save_data(data)
            
            # Определяем расширение
            ext = detect_extension(text, code)
            filename = f"code_{datetime.now().strftime('%H%M%S')}.{ext}"
            
            # Сохраняем код во временный файл
            with open(filename, "w", encoding="utf-8") as f:
                f.write(code)
            
            # Получаем информацию о файле
            file_size = os.path.getsize(filename) / 1024  # в KB
            line_count = len(code.splitlines())
            
            # Отправляем файл
            with open(filename, "rb") as f:
                await app.bot.send_document(
                    chat_id=chat_id,
                    document=f,
                    filename=filename,
                    caption=(
                        f"📁 `{filename}`\n"
                        f"📊 {file_size:.1f} KB • {line_count} строк\n"
                        f"🤖 {ai_name}"
                    )
                )
            
            # Удаляем временный файл
            os.remove(filename)
            
            # Отправляем объяснение, если оно есть
            explanation = re.sub(r'```.*?```', '', result, flags=re.DOTALL).strip()
            if explanation and len(explanation) > 50:
                await send_long_message(
                    app.bot,
                    chat_id,
                    f"✅ *Код готов!*\n\n{explanation}",
                )
        
        else:
            # Отправляем обычный текстовый ответ с проактивными кнопками
            header = f"✅ *Готово!* (_{ai_name}_)\n\n"
            lower_q = text.lower() if text else ""
            suggestions = []
            if any(w in lower_q for w in ['задач', 'сделать нужно', 'план', 'todo']):
                suggestions.append(("➕ В задачи", "tasks_add_from_msg"))
            if any(w in lower_q for w in ['цел', 'хочу достичь', 'мечт']):
                suggestions.append(("🎯 В цели", "goals_add_from_msg"))
            if any(w in lower_q for w in ['напомни', 'не забыть', 'завтра в', 'через час']):
                suggestions.append(("⏰ Напоминание", "reminder_add"))
            
            if suggestions:
                kb = InlineKeyboardMarkup([[
                    InlineKeyboardButton(lbl, callback_data=cb)
                    for lbl, cb in suggestions[:3]
                ]])
                await send_long_message(app.bot, chat_id, header + result, reply_markup=kb)
            else:
                await send_long_message(app.bot, chat_id, header + result)
        
        # Сохраняем данные
        save_data(data)
        
    except Exception as e:
        logger.error(f"Process message error: {e}")
        try:
            await app.bot.send_message(
                chat_id=chat_id,
                text=f"❌ Произошла ошибка: {e}"
            )
        except:
            pass

# =============================================
#  MAIN
# =============================================
# ---- ПОГОДА ----
# Словарь популярных городов СНГ для перевода
CITY_TRANSLATE = {
    "ташкент": "Tashkent", "москва": "Moscow", "питер": "Saint Petersburg",
    "санкт-петербург": "Saint Petersburg", "алматы": "Almaty", "алма-ата": "Almaty",
    "бишкек": "Bishkek", "астана": "Astana", "нур-султан": "Astana",
    "минск": "Minsk", "киев": "Kyiv", "баку": "Baku", "ереван": "Yerevan",
    "тбилиси": "Tbilisi", "душанбе": "Dushanbe", "ашхабад": "Ashgabat",
    "самарканд": "Samarkand", "бухара": "Bukhara", "андижан": "Andijan",
    "новосибирск": "Novosibirsk", "екатеринбург": "Yekaterinburg",
    "казань": "Kazan", "краснодар": "Krasnodar", "сочи": "Sochi",
}

async def get_weather(city: str = "Tashkent") -> str:
    """Получает погоду через OpenWeatherMap"""
    try:
        # Переводим кириллицу на английский
        city_en = CITY_TRANSLATE.get(city.lower().strip(), city)
        url = f"https://api.openweathermap.org/data/2.5/forecast?q={city_en}&appid={WEATHER_KEY}&units=metric&lang=ru&cnt=8"
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get(url)
        if r.status_code != 200:
            # Fallback: пробуем через Geocoding API
            try:
                geo_url = f"https://api.openweathermap.org/geo/1.0/direct?q={city}&limit=1&appid={WEATHER_KEY}"
                async with httpx.AsyncClient(timeout=10) as gc:
                    geo_r = await gc.get(geo_url)
                if geo_r.status_code == 200 and geo_r.json():
                    g = geo_r.json()[0]
                    url = f"https://api.openweathermap.org/data/2.5/forecast?lat={g['lat']}&lon={g['lon']}&appid={WEATHER_KEY}&units=metric&lang=ru&cnt=8"
                    async with httpx.AsyncClient(timeout=10) as gc2:
                        r = await gc2.get(url)
            except Exception as ge:
                logger.error(f"Geocoding fallback error: {ge}")
            if r.status_code != 200:
                return f"❌ Город не найден: {city}\nПопробуй написать по-английски (Tashkent, Moscow)"
        d = r.json()
        city_name = d["city"]["name"]
        country = d["city"]["country"]
        lines = [f"🌤 *Погода в {city_name}, {country}*\n"]
        seen_dates = []
        for item in d["list"]:
            dt = datetime.fromtimestamp(item["dt"])
            date_str = dt.strftime("%d.%m")
            time_str = dt.strftime("%H:%M")
            if date_str not in seen_dates:
                seen_dates.append(date_str)
                lines.append(f"\n📅 *{dt.strftime('%d.%m (%A)')}*")
            temp = round(item["main"]["temp"])
            feels = round(item["main"]["feels_like"])
            desc = item["weather"][0]["description"].capitalize()
            wind = item["wind"]["speed"]
            humidity = item["main"]["humidity"]
            emoji = "☀️" if "ясно" in desc.lower() else "🌧" if "дожд" in desc.lower() else "❄️" if "снег" in desc.lower() else "☁️"
            lines.append(f"  {time_str} {emoji} {temp}°C (ощущ. {feels}°C), {desc}, 💨{wind}м/с, 💧{humidity}%")
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"Weather error: {e}")
        return "❌ Не удалось получить погоду. Попробуй позже."


async def morning_weather_job(app) -> None:
    """Утренняя рассылка погоды в 7:00 по локальному времени пользователя"""
    while True:
        try:
            # Railway сервер работает в UTC — берём offset из данных
            d_tz = load_data()
            tz_offset = d_tz.get("timezone_offset", 5)  # UTC+5 Ташкент по умолчанию
            # Текущее локальное время пользователя
            utc_now = datetime.utcnow()
            local_now = utc_now + timedelta(hours=tz_offset)
            # Ждём до 7:00 по локальному времени
            target = local_now.replace(hour=7, minute=0, second=0, microsecond=0)
            if local_now >= target:
                target += timedelta(days=1)
            wait_secs = (target - local_now).total_seconds()
            logger.info(f"Morning job: UTC={utc_now.strftime('%H:%M')} Local={local_now.strftime('%H:%M')} (UTC+{tz_offset}), sleeping {wait_secs/3600:.1f}h")
            await asyncio.sleep(wait_secs)

            if OWNER_CHAT_ID:
                data = load_data()
                city = data.get("weather_city", "Tashkent")
                weather = await get_weather(city)
                d_tmp = load_data()
                user_name = d_tmp.get("profile", {}).get("name", "Господин")
                active_tasks = [t for t in d_tmp.get("tasks", []) if not t.get("done")]
                greeting = f"🌅 *Доброе утро, {user_name}!*\n\n{weather}\n"
                if active_tasks:
                    greeting += f"\n📋 *Активных задач: {len(active_tasks)}*\n"
                    for t in active_tasks[:3]:
                        greeting += f"• {t['text']}\n"
                greeting += "\n_Хорошего дня!_ 🎯"
                # AI добавляет краткий совет дня
                tip, _ = await get_best_ai_response(
                    "Дай один короткий мотивационный совет на сегодня — одно предложение.",
                    "Ты личный помощник. Отвечай кратко."
                )
                if tip:
                    greeting += f"\n\n💡 *Совет дня:* {tip}"
                await app.bot.send_message(OWNER_CHAT_ID, greeting, parse_mode="Markdown")
        except Exception as e:
            logger.error(f"Morning weather job error: {e}")
            await asyncio.sleep(3600)


# ---- ПЕРЕСКАЗ ССЫЛОК ----
async def fetch_url_content(url: str) -> str:
    """Скачивает и парсит содержимое страницы"""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "ru-RU,ru;q=0.9,en;q=0.8",
        }
        async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
            r = await client.get(url, headers=headers)
        
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(r.text, "html.parser")
            # Убираем ненужное
            for tag in soup(["script", "style", "nav", "footer", "header", "aside", "ads"]):
                tag.decompose()
            # Берём основной текст
            main = soup.find("article") or soup.find("main") or soup.find("body")
            if main:
                text = main.get_text(separator=" ", strip=True)
            else:
                text = soup.get_text(separator=" ", strip=True)
            import re
            text = re.sub(r"\s+", " ", text).strip()
        except ImportError:
            import re
            text = re.sub(r"<script[^>]*>.*?</script>", "", r.text, flags=re.DOTALL)
            text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"\s+", " ", text).strip()
        
        return text[:8000] if text else "error:empty"
    except Exception as e:
        return f"error:{e}"


async def get_youtube_transcript(video_id: str) -> str:
    """Получает субтитры YouTube видео"""
    # Способ 1: youtube-transcript-api (реальные субтитры)
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        # Пробуем русские субтитры, потом английские
        for lang in ['ru', 'en', 'a.ru', 'a.en']:
            try:
                transcript = transcript_list.find_transcript([lang])
                entries = transcript.fetch()
                text = " ".join(e['text'] for e in entries)
                return text[:6000]
            except:
                continue
    except Exception as e:
        logger.warning(f"Transcript API error: {e}")

    # Способ 2: YouTube Data API (название + описание)
    try:
        url = f"https://www.googleapis.com/youtube/v3/videos?part=snippet&id={video_id}&key={YOUTUBE_KEY}"
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get(url)
        d = r.json()
        if d.get("items"):
            item = d["items"][0]
            title = item["snippet"]["title"]
            desc = item["snippet"]["description"][:3000]
            channel = item["snippet"]["channelTitle"]
            return f"Название: {title}\nКанал: {channel}\nОписание: {desc}"
    except Exception as e:
        logger.error(f"YouTube API error: {e}")

    return ""


async def summarize_url(url: str) -> str:
    """Пересказывает содержимое ссылки"""
    import re
    is_youtube = "youtube.com" in url or "youtu.be" in url
    content = ""

    if is_youtube:
        # Извлекаем video_id
        yt_match = re.search(r'(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})', url)
        if yt_match:
            video_id = yt_match.group(1)
            content = await get_youtube_transcript(video_id)
        if not content:
            content = await fetch_url_content(url)
    else:
        content = await fetch_url_content(url)

    if content.startswith("error:"):
        return f"❌ Не удалось открыть ссылку: {content[6:]}"
    if len(content) < 100:
        return "❌ Страница пустая или защищена от парсинга."

    source_type = "YouTube видео" if is_youtube else "страницу"
    prompt = (
        f"Я дам тебе содержимое {source_type}. Сделай подробный пересказ:\n"
        f"1. 📌 О чём материал (2-3 предложения)\n"
        f"2. 🔑 Главные идеи и факты (5-7 пунктов)\n"
        f"3. 💡 Выводы и что важно запомнить\n\n"
        f"Содержимое:\n{content[:5000]}"
    )
    result, ai_name = await get_best_ai_response(prompt, "Ты эксперт по анализу и пересказу контента. Отвечай на русском.")
    if result:
        return f"🔗 *Пересказ материала*\n_(AI: {ai_name})_\n\n{result}"
    return "❌ Не удалось создать пересказ."


# ---- АНАЛИЗ ФАЙЛОВ ----
async def analyze_document(file_bytes: bytes, filename: str, user_prompt: str = "") -> str:
    """Анализирует документ — PDF, TXT, DOCX и др."""
    import io
    text = ""
    fname_lower = filename.lower()

    try:
        if fname_lower.endswith(".txt") or fname_lower.endswith(".md"):
            text = file_bytes.decode("utf-8", errors="ignore")
        elif fname_lower.endswith(".pdf"):
            try:
                import pypdf
                reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                pages_text = []
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        pages_text.append(t)
                text = "\n".join(pages_text)
                if not text.strip():
                    return "❌ PDF защищён или содержит только изображения — текст извлечь не удалось."
            except ImportError:
                return "❌ Библиотека pypdf не установлена."
            except Exception as e:
                return f"❌ Ошибка чтения PDF: {e}"
        elif fname_lower.endswith(".docx"):
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_bytes))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                # Fallback: читаем как zip и извлекаем XML
                try:
                    import zipfile, xml.etree.ElementTree as ET
                    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                        xml_content = z.read("word/document.xml")
                    root = ET.fromstring(xml_content)
                    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
                    paragraphs = []
                    for para in root.iter(f"{ns}p"):
                        texts = [t.text for t in para.iter(f"{ns}t") if t.text]
                        if texts:
                            paragraphs.append("".join(texts))
                    text = "\n".join(paragraphs)
                except Exception as e2:
                    return f"❌ Не удалось прочитать DOCX: {e2}"
            except Exception as e:
                return f"❌ Ошибка чтения DOCX: {e}"
        elif fname_lower.endswith(".csv"):
            text = file_bytes.decode("utf-8", errors="ignore")[:5000]
        elif fname_lower.endswith((".xlsx", ".xls")):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
                rows = []
                for sheet in wb.sheetnames[:3]:  # макс 3 листа
                    ws = wb[sheet]
                    rows.append(f"=== Лист: {sheet} ===")
                    for row in ws.iter_rows(max_row=200, values_only=True):
                        row_text = " | ".join(str(c) if c is not None else "" for c in row)
                        if row_text.strip(" |"):
                            rows.append(row_text)
                text = "\n".join(rows)[:6000]
            except ImportError:
                return "❌ Библиотека openpyxl не установлена."
            except Exception as e:
                return f"❌ Ошибка чтения Excel: {e}"
        elif fname_lower.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text.strip())
                    if slide_content:
                        slides_text.append(f"[Слайд {i}] " + " | ".join(slide_content))
                text = "\n".join(slides_text)[:6000]
            except ImportError:
                return "❌ Библиотека python-pptx не установлена."
            except Exception as e:
                return f"❌ Ошибка чтения PPTX: {e}"
        elif fname_lower.endswith(".json"):
            try:
                import json as _json
                parsed = _json.loads(file_bytes.decode("utf-8", errors="ignore"))
                text = _json.dumps(parsed, ensure_ascii=False, indent=2)[:6000]
            except Exception as e:
                text = file_bytes.decode("utf-8", errors="ignore")[:5000]
        else:
            # Код и другие текстовые форматы
            text = file_bytes.decode("utf-8", errors="ignore")[:6000]
    except Exception as e:
        return f"❌ Ошибка чтения файла: {e}"

    if not text.strip():
        return "❌ Файл пустой или не удалось извлечь текст."

    action = user_prompt.lower() if user_prompt else ""

    if "конспект" in action or "краткое" in action:
        prompt = f"Сделай подробный конспект этого документа с заголовками и ключевыми пунктами:\n\n{text[:6000]}"
    elif "викторин" in action or "тест" in action or "вопрос" in action:
        prompt = f"Создай викторину из 10 вопросов с вариантами ответов (А/Б/В/Г) по этому материалу. В конце дай правильные ответы:\n\n{text[:6000]}"
    elif "перевод" in action or "перевести" in action:
        prompt = f"Переведи этот документ на русский язык:\n\n{text[:6000]}"
    else:
        prompt = (
            f"Проанализируй документ '{filename}':\n"
            f"1. 📌 О чём документ\n"
            f"2. 🔑 Ключевые моменты\n"
            f"3. 📊 Структура и содержание\n"
            f"4. 💡 Выводы\n\n"
            f"Документ:\n{text[:6000]}"
        )

    result, ai_name = await get_best_ai_response(prompt, "Ты эксперт по анализу документов. Отвечай на русском.")
    if result:
        return f"📄 *Анализ файла: {filename}*\n_(AI: {ai_name})_\n\n{result}"
    return "❌ Не удалось проанализировать файл."


# ---- ГЕНЕРАТОР ПРОМПТОВ ----
async def generate_prompt(description: str) -> str:
    """Генерирует оптимальный промпт по описанию"""
    prompt = (
        f"Пользователь хочет создать промпт для AI. Его описание:\n{description}\n\n"
        f"Создай 3 варианта промпта — от простого к сложному:\n"
        f"**Вариант 1 (простой):** ...\n"
        f"**Вариант 2 (средний):** ...\n"
        f"**Вариант 3 (продвинутый):** ...\n\n"
        f"Каждый промпт должен быть готов к использованию, на русском."
    )
    result, ai_name = await get_best_ai_response(prompt, "Ты эксперт по prompt engineering.")
    if result:
        return f"✍️ *Генератор промптов*\n_(AI: {ai_name})_\n\n{result}"
    return "❌ Не удалось создать промпт."


# ---- КОМАНДЫ ----
async def weather_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Команда /weather"""
    args = context.args
    city = " ".join(args) if args else None

    if not city:
        # Проверяем сохранённый город
        data = load_data()
        city = data.get("weather_city", "")

    if not city:
        await update.message.reply_text(
            "🌤 *Погода*\n\nНапиши название города:\nПример: Москва, Алматы, London",
            parse_mode="Markdown"
        )
        context.user_data["waiting_for"] = "weather_city"
        return

    msg = await update.message.reply_text(f"🌤 Получаю погоду для *{city}*...", parse_mode="Markdown")
    weather = await get_weather(city)

    # Сохраняем город
    data = load_data()
    data["weather_city"] = city
    save_data(data)

    await msg.edit_text(weather, parse_mode="Markdown")


async def summarize_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Команда /summarize или кнопка Пересказ ссылки"""
    await update.message.reply_text(
        "🔗 *Пересказ ссылки*\n\nОтправь ссылку на:\n"
        "• YouTube видео\n• Статью или сайт\n• Новость\n\nПросто вставь URL:",
        parse_mode="Markdown"
    )
    context.user_data["waiting_for"] = "summarize_url"


async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик файлов — анализ документов"""
    doc = update.message.document
    if not doc:
        return

    allowed = [".txt", ".pdf", ".docx", ".md", ".csv", ".xlsx", ".xls", ".pptx", ".json", ".py", ".js", ".html", ".xml"]
    fname = doc.file_name or "document"
    ext = "." + fname.rsplit(".", 1)[-1].lower() if "." in fname else ""

    if ext not in allowed:
        await update.message.reply_text(
            f"❌ Формат {ext} не поддерживается.\n"
            f"Поддерживаю: PDF, DOCX, XLSX, PPTX, TXT, CSV, JSON, MD, PY, JS, HTML, XML"
        )
        return

    if doc.file_size > 10 * 1024 * 1024:  # 10MB
        await update.message.reply_text("❌ Файл слишком большой (макс. 10MB)")
        return

    msg = await update.message.reply_text(f"📄 Читаю *{fname}*...", parse_mode="Markdown")

    try:
        file = await context.bot.get_file(doc.file_id)
        file_bytes = await file.download_as_bytearray()

        # Берём caption как подсказку что делать
        user_prompt = update.message.caption or ""

        # Для данных — график или анализ
        if ext in [".xlsx", ".xls", ".csv"] and ("график" in user_prompt.lower() or not user_prompt):
            if "график" in user_prompt.lower() or not user_prompt:
                chart_path, chart_err = await create_chart_from_data(bytes(file_bytes), fname)
                if chart_path:
                    with open(chart_path, 'rb') as cf:
                        await context.bot.send_photo(
                            chat_id=update.effective_chat.id, photo=cf,
                            caption=f"📈 *График: {fname}*", parse_mode="Markdown"
                        )
                    import os; os.remove(chart_path)
                    await msg.delete()
                    if not user_prompt:
                        result = await analyze_data_file(bytes(file_bytes), fname)
                        await send_long_message(context.application.bot, update.effective_chat.id, result, parse_mode="Markdown")
                    return
                else:
                    result = chart_err or await analyze_data_file(bytes(file_bytes), fname)
            else:
                result = await analyze_data_file(bytes(file_bytes), fname)
        else:
            result = await analyze_document(bytes(file_bytes), fname, user_prompt)
        data = load_data()
        update_stats(data, "file")
        save_data(data)
        await msg.delete()
        await send_long_message(context.application.bot, update.effective_chat.id, result, parse_mode="Markdown")

    except Exception as e:
        logger.error(f"Document handler error: {e}")
        await msg.edit_text(f"❌ Ошибка: {e}")



# ---- СТАТИСТИКА: обновление ----
def update_stats(data: dict, msg_type: str = "message") -> None:
    """Обновляет статистику использования"""
    if "stats" not in data:
        data["stats"] = {
            "messages_total": 0, "voice_total": 0, "files_total": 0,
            "links_total": 0, "topics_studied": 0, "tasks_done": 0,
            "days_active": 0, "first_seen": datetime.now().strftime("%d.%m.%Y"),
            "last_seen": datetime.now().strftime("%d.%m.%Y"),
            "hourly": {}, "daily_messages": {}
        }
    s = data["stats"]
    s["messages_total"] = s.get("messages_total", 0) + 1
    if msg_type == "voice": s["voice_total"] = s.get("voice_total", 0) + 1
    elif msg_type == "file": s["files_total"] = s.get("files_total", 0) + 1
    elif msg_type == "link": s["links_total"] = s.get("links_total", 0) + 1

    now = datetime.now()
    hour = str(now.hour)
    today = now.strftime("%d.%m.%Y")
    s["hourly"][hour] = s["hourly"].get(hour, 0) + 1
    s["daily_messages"][today] = s["daily_messages"].get(today, 0) + 1
    s["last_seen"] = today


def analyze_habits(data: dict) -> str:
    """Анализирует привычки пользователя по статистике"""
    s = data.get("stats", {})
    hourly = s.get("hourly", {})
    daily = s.get("daily_messages", {})

    if not hourly:
        return "📊 Пока недостаточно данных для анализа. Используй бота несколько дней!"

    # Самый активный час
    peak_hour = max(hourly, key=lambda h: hourly[h], default="?")
    peak_count = hourly.get(peak_hour, 0)

    # Среднее сообщений в день
    if daily:
        avg_per_day = sum(daily.values()) / len(daily)
        max_day = max(daily, key=lambda d: daily[d], default="?")
        max_day_count = daily.get(max_day, 0)
    else:
        avg_per_day = 0
        max_day = "?"
        max_day_count = 0

    # Активность по времени суток
    morning = sum(hourly.get(str(h), 0) for h in range(6, 12))
    afternoon = sum(hourly.get(str(h), 0) for h in range(12, 18))
    evening = sum(hourly.get(str(h), 0) for h in range(18, 24))
    night = sum(hourly.get(str(h), 0) for h in range(0, 6))
    peak_period = max([("🌅 Утро (6-12)", morning), ("☀️ День (12-18)", afternoon),
                       ("🌆 Вечер (18-24)", evening), ("🌙 Ночь (0-6)", night)],
                      key=lambda x: x[1])

    topics = data.get("topics", [])
    tasks = data.get("tasks", [])
    done_tasks = len([t for t in tasks if t.get("done")])

    lines = [
        "🧬 *Анализ твоих привычек*\n",
        f"⏰ Самое активное время: *{peak_hour}:00* ({peak_count} сообщений)",
        f"📅 Самый активный день: *{max_day}* ({max_day_count} сообщений)",
        f"📈 В среднем в день: *{avg_per_day:.1f}* сообщений",
        f"🕐 Пик активности: *{peak_period[0]}*",
        f"\n📚 Изучено тем: *{len(topics)}*",
        f"✅ Выполнено задач: *{done_tasks}* из *{len(tasks)}*",
        f"🎤 Голосовых: *{s.get('voice_total', 0)}*",
        f"📄 Файлов: *{s.get('files_total', 0)}*",
        f"🔗 Ссылок: *{s.get('links_total', 0)}*",
    ]

    # Рекомендации
    lines.append("\n💡 *Рекомендации:*")
    if peak_period[1] == evening:
        lines.append("• Ты вечерний человек — планируй важные задачи на вечер")
    elif peak_period[1] == morning:
        lines.append("• Ты жаворонок — используй утро для сложных задач")
    if avg_per_day > 20:
        lines.append("• Высокая активность — молодец, используешь бота по максимуму!")
    if len(topics) > 5:
        lines.append(f"• Уже изучил {len(topics)} тем — продолжай самообразование!")

    return "\n".join(lines)


# ---- ПРОФИЛЬ ----
async def profile_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Показывает и управляет профилем"""
    data = load_data()
    user = update.effective_user
    profile = data.get("profile", {})
    stats = data.get("stats", {})

    # Берём имя из Telegram
    tg_name = user.first_name or user.username or "Господин"
    if user.last_name:
        tg_name += f" {user.last_name}"

    # Обновляем имя в профиле если не задано
    if not profile.get("name"):
        profile["name"] = tg_name
        data["profile"] = profile
        save_data(data)

    name = profile.get("name", tg_name)
    city = profile.get("city") or data.get("weather_city", "не указан")
    interests = profile.get("interests", [])
    goals = profile.get("goals", [])
    tasks = data.get("tasks", [])
    done_tasks = len([t for t in tasks if t.get("done")])

    text = (
        f"👤 *Личный профиль*\n\n"
        f"👋 Имя: *{name}*\n"
        f"🌍 Город: *{city}*\n"
        f"📅 В боте с: *{stats.get('first_seen', 'сегодня')}*\n"
        f"📊 Сообщений всего: *{stats.get('messages_total', 0)}*\n"
        f"✅ Задач выполнено: *{done_tasks}/{len(tasks)}*\n"
        f"📚 Тем изучено: *{len(data.get('topics', []))}*\n"
    )
    if interests:
        text += f"\n🎯 Интересы: {', '.join(interests[:5])}"
    if goals:
        text += f"\n🚀 Цели: {', '.join(goals[:3])}"

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("✏️ Изменить имя", callback_data="profile_name"),
         InlineKeyboardButton("🌍 Изменить город", callback_data="profile_city")],
        [InlineKeyboardButton("🎯 Добавить интерес", callback_data="profile_interest"),
         InlineKeyboardButton("🚀 Добавить цель", callback_data="profile_goal")],
        [InlineKeyboardButton("🧬 Анализ привычек", callback_data="profile_habits")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


# ---- ЗАДАЧИ ----
async def tasks_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Показывает список задач"""
    data = load_data()
    tasks = data.get("tasks", [])

    if not tasks:
        text = (
            "✅ *Ежедневные задачи*\n\n"
            "У тебя пока нет задач!\n\n"
            "Добавь задачу:\n"
            "• Напиши: _добавь задачу [название]_\n"
            "• Или: _задача [название]_"
        )
        keyboard = InlineKeyboardMarkup([
            [InlineKeyboardButton("➕ Добавить задачу", callback_data="task_add")]
        ])
    else:
        active = [t for t in tasks if not t.get("done")]
        done = [t for t in tasks if t.get("done")]

        lines = ["✅ *Мои задачи*\n"]
        if active:
            lines.append("📋 *Активные:*")
            for i, t in enumerate(active, 1):
                priority = "🔴" if t.get("priority") == "high" else "🟡" if t.get("priority") == "medium" else "🟢"
                lines.append(f"{priority} {i}. {t['text']}")
        if done:
            lines.append(f"\n✅ *Выполнено ({len(done)}):*")
            for t in done[-3:]:  # последние 3
                lines.append(f"~~{t['text']}~~")

        text = "\n".join(lines)
        keyboard = InlineKeyboardMarkup([
            [InlineKeyboardButton("➕ Добавить", callback_data="task_add"),
             InlineKeyboardButton("✅ Отметить выполненной", callback_data="task_done")],
            [InlineKeyboardButton("🗑 Очистить выполненные", callback_data="task_clear")],
        ])

    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


# ---- СТАТИСТИКА ----
async def stats_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Показывает статистику"""
    data = load_data()
    s = data.get("stats", {})
    tasks = data.get("tasks", [])
    done_tasks = len([t for t in tasks if t.get("done")])

    # Активность сегодня
    today = datetime.now().strftime("%d.%m.%Y")
    today_msgs = s.get("daily_messages", {}).get(today, 0)

    # Самый активный час
    hourly = s.get("hourly", {})
    peak = max(hourly, key=lambda h: hourly[h], default="—") if hourly else "—"

    text = (
        "📊 *Статистика использования*\n\n"
        f"📅 Сегодня сообщений: *{today_msgs}*\n"
        f"📨 Всего сообщений: *{s.get('messages_total', 0)}*\n"
        f"🎤 Голосовых: *{s.get('voice_total', 0)}*\n"
        f"📄 Файлов проанализировано: *{s.get('files_total', 0)}*\n"
        f"🔗 Ссылок обработано: *{s.get('links_total', 0)}*\n"
        f"📚 Тем изучено: *{len(data.get('topics', []))}*\n"
        f"✅ Задач выполнено: *{done_tasks}*\n"
        f"⏰ Пик активности: *{peak}:00*\n"
        f"📆 В боте с: *{s.get('first_seen', today)}*\n"
        f"🕐 Последний визит: *{s.get('last_seen', today)}*\n"
    )

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("🧬 Анализ привычек", callback_data="profile_habits")]
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)



# ============================================================
# ФИЧА 1: КОНТЕКСТ МЕЖДУ СЕССИЯМИ + СУПЕРПАМЯТЬ
# ============================================================

def extract_facts_from_message(text: str, data: dict) -> None:
    """Извлекает факты из сообщения и сохраняет в долгосрочную память"""
    brain = data.setdefault("second_brain", {
        "facts": [], "knowledge": {}, "context_log": [], "goals_progress": {}
    })
    
    # Логируем контекст
    ctx = brain.setdefault("context_log", [])
    ctx.append({"text": text[:200], "time": datetime.now().strftime("%d.%m.%Y %H:%M")})
    if len(ctx) > 200:
        brain["context_log"] = ctx[-200:]
    
    # Автоматически извлекаем факты по паттернам
    facts = brain.setdefault("facts", [])
    t = text.lower()
    
    patterns = [
        # Имя
        (r"меня зовут ([а-яёa-z]+)", lambda m: f"Имя пользователя: {m.group(1).title()}"),
        (r"я ([а-яё]+(?:ин|ов|ев|ян|ский))", lambda m: f"Возможное имя: {m.group(1).title()}"),
        # Возраст
        (r"мне (\d{1,2}) ле", lambda m: f"Возраст: {m.group(1)} лет"),
        # Город
        (r"живу в ([а-яё]+)", lambda m: f"Город: {m.group(1).title()}"),
        (r"нахожусь в ([а-яё]+)", lambda m: f"Местоположение: {m.group(1).title()}"),
        # Работа
        (r"работаю (?:в |как |)([а-яё\w\s]{3,25})", lambda m: f"Работа: {m.group(1).strip()}"),
        (r"я (?:программист|разработчик|дизайнер|менеджер|врач|учитель|студент)", 
         lambda m: f"Профессия: {m.group(0).replace('я ','').title()}"),
        # Предпочтения
        (r"я люблю ([а-яё\w\s]{3,30})", lambda m: f"Любит: {m.group(1).strip()}"),
        (r"не люблю ([а-яё\w\s]{3,30})", lambda m: f"Не любит: {m.group(1).strip()}"),
        # Команда запомни
        (r"запомни[,:]?\s+(.{5,100})", lambda m: m.group(1).strip()),
        (r"важно[,:]?\s+(.{5,100})", lambda m: f"Важно: {m.group(1).strip()}"),
    ]
    
    import re
    for pattern, extractor in patterns:
        match = re.search(pattern, t)
        if match:
            try:
                fact = extractor(match)
                # Не дублируем факты по ключевому слову
                key = fact.split(":")[0] if ":" in fact else fact[:20]
                facts = [f for f in facts if not f.lower().startswith(key.lower())]
                facts.append(fact)
            except:
                pass
    
    if len(facts) > 50:
        brain["facts"] = facts[-50:]
    else:
        brain["facts"] = facts


async def build_session_context(data: dict) -> str:
    """Строит контекст из прошлых сессий для AI"""
    brain = data.get("second_brain", {})
    ctx = brain.get("context_log", [])
    facts = brain.get("facts", [])
    goals = data.get("goals_system", {}).get("goals", [])

    lines = []
    if facts:
        lines.append("Известные факты о пользователе: " + "; ".join(facts[-10:]))
    if goals:
        active = [g for g in goals if g.get("status") == "active"]
        if active:
            lines.append("Текущие цели: " + "; ".join(g["title"] for g in active[:3]))
    if ctx:
        recent = ctx[-5:]
        lines.append("Недавние темы разговора: " + "; ".join(c["text"][:50] for c in recent))
    return "\n".join(lines)


# ============================================================
# ФИЧА 2: SECOND BRAIN — СИСТЕМА ЗНАНИЙ
# ============================================================

async def save_to_brain(data: dict, topic: str, content: str) -> None:
    """Сохраняет знание в Second Brain"""
    brain = data.setdefault("second_brain", {"facts": [], "knowledge": {}, "context_log": [], "goals_progress": {}})
    knowledge = brain.setdefault("knowledge", {})
    topic_key = topic.lower().strip()
    if topic_key not in knowledge:
        knowledge[topic_key] = {
            "title": topic,
            "notes": [],
            "created": datetime.now().strftime("%d.%m.%Y")
        }
    knowledge[topic_key]["notes"].append({
        "text": content[:500],
        "date": datetime.now().strftime("%d.%m.%Y %H:%M")
    })
    knowledge[topic_key]["updated"] = datetime.now().strftime("%d.%m.%Y")


async def brain_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Second Brain — система знаний"""
    data = load_data()
    brain = data.get("second_brain", {})
    knowledge = brain.get("knowledge", {})
    facts = brain.get("facts", [])

    if not knowledge and not facts:
        text = (
            "📚 *Second Brain — Система знаний*\n\n"
            "Здесь хранятся все твои знания!\n\n"
            "Как пополнять:\n"
            "• Изучай темы через _изучи [тема]_\n"
            "• Анализируй файлы и ссылки\n"
            "• Пиши _запомни: [факт]_ для сохранения\n\n"
            "База знаний пока пуста — начни изучать!"
        )
    else:
        topics_list = list(knowledge.keys())[:15]
        lines = [f"📚 *Second Brain*\n", f"🗂 Тем в базе: *{len(knowledge)}*", f"💡 Фактов: *{len(facts)}*\n"]
        if topics_list:
            lines.append("*Последние темы:*")
            for t in topics_list[-10:]:
                k = knowledge[t]
                lines.append(f"  • {k['title']} ({len(k['notes'])} заметок)")
        text = "\n".join(lines)

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("🔍 Найти в базе знаний", callback_data="brain_search"),
         InlineKeyboardButton("💡 Все факты", callback_data="brain_facts")],
        [InlineKeyboardButton("📖 Список тем", callback_data="brain_topics")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


# ============================================================
# ФИЧА 3: СИСТЕМА ЦЕЛЕЙ С ТРЕКИНГОМ
# ============================================================

async def goals_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Система целей с трекингом"""
    data = load_data()
    gs = data.get("goals_system", {"goals": []})
    goals = gs.get("goals", [])
    active = [g for g in goals if g.get("status") == "active"]
    done = [g for g in goals if g.get("status") == "done"]

    if not goals:
        text = (
            "🎯 *Система целей*\n\n"
            "Постав цель и я буду отслеживать прогресс!\n\n"
            "Примеры целей:\n"
            "• Выучить Python за 3 месяца\n"
            "• Читать 1 книгу в месяц\n"
            "• Заниматься спортом 3 раза в неделю"
        )
    else:
        lines = [f"🎯 *Мои цели*\n"]
        if active:
            lines.append(f"🔥 *Активные ({len(active)}):*")
            for g in active:
                progress = g.get("progress", 0)
                bar = "█" * (progress // 10) + "░" * (10 - progress // 10)
                deadline = g.get("deadline", "без срока")
                lines.append(f"\n*{g['title']}*")
                lines.append(f"[{bar}] {progress}%")
                lines.append(f"📅 Срок: {deadline}")
                if g.get("steps"):
                    done_steps = len([s for s in g["steps"] if s.get("done")])
                    lines.append(f"✅ Шагов: {done_steps}/{len(g['steps'])}")
        if done:
            lines.append(f"\n✅ *Выполнено: {len(done)}*")
        text = "\n".join(lines)

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("➕ Новая цель", callback_data="goal_add"),
         InlineKeyboardButton("📈 Обновить прогресс", callback_data="goal_progress")],
        [InlineKeyboardButton("🤖 Автопилот цели", callback_data="goal_autopilot")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


# ============================================================
# ФИЧА 4: ГЛУБОКИЙ ИССЛЕДОВАТЕЛЬ
# ============================================================

async def deep_research(app, chat_id: int, topic: str, data: dict) -> None:
    """Глубокое исследование темы с анализом источников"""
    try:
        await app.bot.send_message(chat_id,
            f"🔬 *Глубокое исследование:* _{topic}_\n\n"
            f"Анализирую 8+ источников...",
            parse_mode="Markdown"
        )

        # Генерируем разнообразные запросы
        queries_prompt = (
            f"Тема: {topic}\n"
            f"Дай 6 поисковых запросов с разных углов:\n"
            f"1. Основное определение\n2. Как работает/механизм\n"
            f"3. Практическое применение\n4. Критика и проблемы\n"
            f"5. Последние исследования\n6. Сравнение с альтернативами\n"
            f"Только список запросов, каждый с новой строки."
        )
        queries_raw, _ = await get_best_ai_response(queries_prompt, "Ты исследователь. Только список.")
        queries = [q.strip() for q in (queries_raw or topic).split("\n") if q.strip()][:6]

        # Собираем данные
        all_results = []
        for i, q in enumerate(queries, 1):
            await app.bot.send_message(chat_id, f"🔍 Источник {i}/{len(queries)}: _{q}_", parse_mode="Markdown")
            result = await search_web(q)
            if result and "ничего не найдено" not in result:
                all_results.append(f"[{q}]\n{result}")
            await asyncio.sleep(0.5)

        if not all_results:
            await app.bot.send_message(chat_id, "❌ Не удалось найти достаточно информации.")
            return

        combined = "\n\n".join(all_results[:5])
        await app.bot.send_message(chat_id, "🧠 Синтезирую аналитический отчёт...")

        research_prompt = (
            f"Тема исследования: {topic}\n\n"
            f"Данные из {len(all_results)} источников:\n{combined[:6000]}\n\n"
            f"Напиши АНАЛИТИЧЕСКИЙ ОТЧЁТ:\n"
            f"## 1. РЕЗЮМЕ (3-4 предложения)\n"
            f"## 2. ДЕТАЛЬНЫЙ АНАЛИЗ\n"
            f"## 3. КЛЮЧЕВЫЕ ФАКТЫ И ЦИФРЫ\n"
            f"## 4. РАЗНЫЕ ТОЧКИ ЗРЕНИЯ (за и против)\n"
            f"## 5. ПРАКТИЧЕСКОЕ ПРИМЕНЕНИЕ\n"
            f"## 6. ВЫВОДЫ И РЕКОМЕНДАЦИИ\n"
            f"## 7. ИСТОЧНИКИ ДЛЯ УГЛУБЛЁННОГО ИЗУЧЕНИЯ\n\n"
            f"Будь аналитичным, указывай на противоречия между источниками."
        )
        report, ai_name = await get_best_ai_response(
            research_prompt,
            "Ты академический исследователь и аналитик. Пиши структурированно и критически."
        )

        if not report:
            await app.bot.send_message(chat_id, "❌ Не удалось создать отчёт.")
            return

        # Сохраняем в Second Brain
        await save_to_brain(data, topic, report[:500])
        save_data(data)

        # Отправляем файлом
        ts = datetime.now().strftime("%d%m%Y_%H%M")
        fname = f"research_{ts}.txt"
        file_content = (
            f"АНАЛИТИЧЕСКИЙ ОТЧЁТ: {topic}\n{'='*60}\n"
            f"Дата: {datetime.now().strftime('%d.%m.%Y %H:%M')} | AI: {ai_name}\n"
            f"Источников: {len(all_results)}\n{'='*60}\n\n{report}\n\n"
            f"{'='*60}\nРАВ ДАННЫЕ\n{'='*60}\n{combined[:3000]}"
        )
        with open(fname, "w", encoding="utf-8") as f:
            f.write(file_content)
        with open(fname, "rb") as f:
            await app.bot.send_document(
                chat_id=chat_id, document=f,
                filename=f"Исследование_{topic[:25]}.txt",
                caption=f"🔬 *Исследование завершено: {topic}*\n🤖 {ai_name}\n📊 Источников: {len(all_results)}",
                parse_mode="Markdown"
            )
        import os
        os.remove(fname)

        # Краткое резюме в чат
        summary = report[:600] + ("..." if len(report) > 600 else "")
        await app.bot.send_message(chat_id, f"📌 *Резюме:*\n\n{summary}", parse_mode="Markdown")

    except Exception as e:
        logger.error(f"deep_research error: {e}")
        await app.bot.send_message(chat_id, f"❌ Ошибка исследования: {e}")


def is_research_request(text: str):
    """Определяет запрос на глубокое исследование"""
    patterns = [r'исследуй\s+(.+)', r'глубокий анализ\s+(.+)', r'подробно о\s+(.+)',
                r'аналитика\s+(.+)', r'исследование\s+(.+)']
    import re
    for p in patterns:
        m = re.search(p, text.lower())
        if m:
            return True, m.group(1).strip()
    return False, ""


# ============================================================
# ФИЧА 5: АНАЛИТИК ДАННЫХ (Excel/CSV)
# ============================================================

async def analyze_data_file(file_bytes: bytes, filename: str) -> str:
    """Глубокий анализ данных из Excel/CSV"""
    import io
    fname_lower = filename.lower()
    rows_data = []

    try:
        if fname_lower.endswith(".csv"):
            import csv
            text = file_bytes.decode("utf-8", errors="ignore")
            reader = csv.DictReader(io.StringIO(text))
            for i, row in enumerate(reader):
                if i >= 500: break
                rows_data.append(row)
        elif fname_lower.endswith((".xlsx", ".xls")):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            ws = wb.active
            headers = None
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i == 0:
                    headers = [str(c) if c else f"Col{j}" for j, c in enumerate(row)]
                    continue
                if i > 500: break
                if headers and any(c is not None for c in row):
                    rows_data.append(dict(zip(headers, row)))
    except Exception as e:
        return f"❌ Ошибка чтения файла: {e}"

    if not rows_data:
        return "❌ Файл пустой или не удалось прочитать данные."

    # Базовая статистика
    columns = list(rows_data[0].keys()) if rows_data else []
    total_rows = len(rows_data)

    # Числовые колонки
    numeric_stats = {}
    for col in columns:
        values = []
        for row in rows_data:
            try:
                v = float(str(row.get(col, "")).replace(",", ".").replace(" ", ""))
                values.append(v)
            except: pass
        if len(values) > total_rows * 0.3:
            numeric_stats[col] = {
                "min": min(values), "max": max(values),
                "avg": sum(values)/len(values), "sum": sum(values), "count": len(values)
            }

    # Формируем описание для AI
    stats_text = f"Файл: {filename}\nСтрок: {total_rows}\nКолонки: {', '.join(columns)}\n\n"
    for col, s in numeric_stats.items():
        stats_text += f"{col}: мин={s['min']:.2f}, макс={s['max']:.2f}, среднее={s['avg']:.2f}, сумма={s['sum']:.2f}\n"

    # Первые 10 строк для контекста
    sample = "\n".join(str(r) for r in rows_data[:10])
    stats_text += f"\nПервые строки:\n{sample}"

    prompt = (
        f"Ты аналитик данных. Проанализируй этот датасет:\n\n{stats_text[:4000]}\n\n"
        f"Дай:\n1. 📊 Что представляют эти данные\n"
        f"2. 📈 Ключевые метрики и тренды\n"
        f"3. 🔍 Аномалии и интересные паттерны\n"
        f"4. 💡 Выводы и рекомендации\n"
        f"5. ❓ Вопросы для дальнейшего анализа"
    )
    result, ai_name = await get_best_ai_response(prompt, "Ты опытный Data Analyst. Будь конкретным с цифрами.")
    if result:
        return f"📊 *Анализ данных: {filename}*\n_{ai_name}_\n\n{result}"
    return "❌ Не удалось проанализировать данные."


# ============================================================
# ФИЧА 6: ПЕРСОНАЛЬНЫЙ РЕПЕТИТОР
# ============================================================

async def tutor_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Персональный репетитор"""
    data = load_data()
    courses = data.get("tutor_courses", [])

    if not courses:
        text = (
            "🎓 *Персональный репетитор*\n\n"
            "Я создам персональный курс по любой теме!\n\n"
            "Как работает:\n"
            "1️⃣ Ты называешь тему\n"
            "2️⃣ Я создаю программу курса\n"
            "3️⃣ Объясняю теорию с примерами\n"
            "4️⃣ Даю тест для проверки\n"
            "5️⃣ Напоминаю повторить через 3 дня\n\n"
            "Напиши тему для изучения:"
        )
        context.user_data["waiting_for"] = "tutor_topic"
    else:
        active = [c for c in courses if c.get("status") == "active"]
        lines = ["🎓 *Мои курсы*\n"]
        for c in active[:5]:
            progress = c.get("progress", 0)
            bar = "█" * (progress // 20) + "░" * (5 - progress // 20)
            lines.append(f"📖 *{c['topic']}* [{bar}] {progress}%")
            lines.append(f"   Урок {c.get('current_lesson', 1)}/{c.get('total_lessons', 5)}\n")
        text = "\n".join(lines)

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("➕ Новый курс", callback_data="tutor_new"),
         InlineKeyboardButton("📖 Продолжить урок", callback_data="tutor_continue")],
        [InlineKeyboardButton("📝 Тест", callback_data="tutor_test")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


async def create_tutor_course(app, chat_id: int, topic: str, data: dict) -> None:
    """Создаёт персональный курс"""
    try:
        await app.bot.send_message(chat_id, f"🎓 Создаю курс: *{topic}*...", parse_mode="Markdown")

        prompt = (
            f"Создай план курса по теме: {topic}\n\n"
            f"Структура (5 уроков):\n"
            f"Урок 1: [название] — [краткое описание]\n"
            f"Урок 2: ...\n...\n"
            f"Урок 5: ...\n\n"
            f"Потом напиши Урок 1 полностью:\n"
            f"- Теория (понятно и с примерами)\n"
            f"- 3 практических задания\n"
            f"- Ключевые термины"
        )
        lesson, ai_name = await get_best_ai_response(
            prompt, "Ты опытный преподаватель. Объясняй просто и с примерами."
        )

        if not lesson:
            await app.bot.send_message(chat_id, "❌ Не удалось создать курс.")
            return

        # Сохраняем курс
        course = {
            "topic": topic,
            "status": "active",
            "progress": 0,
            "current_lesson": 1,
            "total_lessons": 5,
            "created": datetime.now().strftime("%d.%m.%Y"),
            "lessons": [{"number": 1, "content": lesson, "done": False}]
        }
        data.setdefault("tutor_courses", []).append(course)
        save_data(data)

        await send_long_message(
            app.bot, chat_id,
            f"🎓 *Курс создан: {topic}*\n\n{lesson}",
            parse_mode="Markdown"
        )
        await app.bot.send_message(
            chat_id,
            "✅ Урок 1 готов!\n\nНапиши /tutor чтобы пройти тест или продолжить обучение.",
        )
    except Exception as e:
        logger.error(f"create_tutor_course error: {e}")
        await app.bot.send_message(chat_id, f"❌ Ошибка: {e}")


# ============================================================
# ФИЧА 7: АВТОПИЛОТ ЗАДАЧ
# ============================================================

async def autopilot_task(app, chat_id: int, task_desc: str, data: dict) -> None:
    """Разбивает задачу на шаги и создаёт план"""
    try:
        await app.bot.send_message(chat_id, f"⚡ *Автопилот:* _{task_desc}_\n\nРазбиваю на шаги...", parse_mode="Markdown")

        prompt = (
            f"Задача: {task_desc}\n\n"
            f"Создай детальный план выполнения:\n"
            f"1. Разбей на 5-7 конкретных шагов\n"
            f"2. Для каждого шага укажи примерное время\n"
            f"3. Укажи что нужно для выполнения\n"
            f"4. Предупреди о возможных проблемах\n\n"
            f"Формат:\n"
            f"Шаг 1: [действие] (⏱ X минут)\n"
            f"Шаг 2: ...\n\n"
            f"Итого: ~X часов/минут"
        )
        plan, ai_name = await get_best_ai_response(
            prompt, "Ты менеджер проектов. Давай конкретные и выполнимые шаги."
        )

        if not plan:
            await app.bot.send_message(chat_id, "❌ Не удалось создать план.")
            return

        # Парсим шаги и добавляем как задачи
        import re
        steps = re.findall(r'Шаг \d+:?\s*(.+?)(?:\(⏱.*?\))?(?:\n|$)', plan)
        tasks = data.setdefault("tasks", [])
        added = 0
        for step in steps[:7]:
            step = step.strip()
            if step and len(step) > 3:
                tasks.append({
                    "text": f"[{task_desc[:20]}] {step}",
                    "done": False,
                    "created": datetime.now().strftime("%d.%m.%Y %H:%M"),
                    "priority": "high",
                    "autopilot": True
                })
                added += 1
        save_data(data)

        await send_long_message(
            app.bot, chat_id,
            f"⚡ *План задачи: {task_desc}*\n_{ai_name}_\n\n{plan}",
            parse_mode="Markdown"
        )
        if added:
            await app.bot.send_message(
                chat_id,
                f"✅ {added} шагов добавлено в список задач!\nПосмотри через /tasks"
            )
    except Exception as e:
        logger.error(f"autopilot_task error: {e}")
        await app.bot.send_message(chat_id, f"❌ Ошибка: {e}")


def is_autopilot_request(text: str):
    """Определяет запрос на автопилот"""
    import re
    for p in [r'автопилот\s+(.+)', r'разбей на шаги\s+(.+)', r'план для\s+(.+)', r'как выполнить\s+(.+)']:
        m = re.search(p, text.lower())
        if m:
            return True, m.group(1).strip()
    return False, ""


# ============================================================
# ФИЧА 1: ДОЛГОСРОЧНАЯ ПАМЯТЬ — автоматически из разговоров
# ============================================================

async def auto_extract_memory(text: str, response: str, data: dict) -> None:
    """Автоматически извлекает факты о пользователе из разговора"""
    try:
        prompt = (
            f"Пользователь написал: {text}\n"
            f"AI ответил: {response[:200]}\n\n"
            f"Извлеки ТОЛЬКО конкретные факты о пользователе (имя, город, работа, интересы, цели, предпочтения).\n"
            f"Если нет фактов — ответь: NONE\n"
            f"Если есть — напиши одной строкой: ФАКТ: [факт]\n"
            f"Максимум 1-2 факта."
        )
        result, _ = await get_best_ai_response(prompt, "Извлекай только реальные факты. Будь краток.")
        if not result or "NONE" in result.upper():
            return
        import re
        facts = re.findall(r'ФАКТ:\s*(.+)', result)
        if facts:
            brain = data.setdefault("second_brain", {"facts": [], "knowledge": {}, "context_log": []})
            existing = brain.setdefault("facts", [])
            for fact in facts:
                fact = fact.strip()
                # Не дублируем
                if fact and fact not in existing and len(fact) > 5:
                    existing.append(fact)
                    if len(existing) > 200:
                        brain["facts"] = existing[-200:]
            logger.info(f"Auto-memory: extracted {len(facts)} facts")
    except Exception as e:
        logger.debug(f"auto_extract_memory error: {e}")


def build_memory_context(data: dict) -> str:
    """Строит расширенный контекст с долгосрочной памятью"""
    brain = data.get("second_brain", {})
    facts = brain.get("facts", [])
    ctx_log = brain.get("context_log", [])
    profile = data.get("profile", {})
    goals = data.get("goals_system", {}).get("goals", [])
    tasks = data.get("tasks", [])

    parts = []

    # Имя пользователя
    name = profile.get("name", "")
    if name:
        parts.append(f"Имя пользователя: {name}")

    # Факты из памяти
    if facts:
        parts.append("Долгосрочная память о пользователе:\n" + "\n".join(f"• {f}" for f in facts[-15:]))

    # Активные цели
    active_goals = [g for g in goals if g.get("status") == "active"]
    if active_goals:
        parts.append("Текущие цели: " + ", ".join(g["title"] for g in active_goals[:3]))

    # Активные задачи
    active_tasks = [t for t in tasks if not t.get("done")][:3]
    if active_tasks:
        parts.append("Активные задачи: " + ", ".join(t["text"][:30] for t in active_tasks))

    # Последние темы из контекста
    if ctx_log:
        recent = ctx_log[-3:]
        parts.append("Недавние темы: " + ", ".join(c["text"][:40] for c in recent))

    return "\n".join(parts)


# ============================================================
# ФИЧА 2: МУЛЬТИАГЕНТНОСТЬ — параллельный запрос к нескольким AI
# ============================================================

async def multi_agent_response(prompt: str, system_context: str) -> Tuple[str, str]:
    """Запрашивает несколько AI параллельно и синтезирует лучший ответ"""
    # Запускаем всех параллельно
    tasks_ai = [
        asyncio.create_task(ask_groq(prompt, system_context)),
        asyncio.create_task(ask_gemini(prompt, system_context)),
        asyncio.create_task(ask_deepseek(prompt, system_context)),
        asyncio.create_task(ask_cohere(prompt, system_context)),
    ]
    names = ["Groq", "Gemini", "DeepSeek", "Cohere"]

    results = await asyncio.gather(*tasks_ai, return_exceptions=True)

    valid = []
    for name, res in zip(names, results):
        if isinstance(res, str) and res and len(res) > 20:
            valid.append((name, res))

    if not valid:
        return None, None

    if len(valid) == 1:
        return valid[0][1], valid[0][0]

    # Синтез лучшего ответа из нескольких
    combined = "\n\n".join(f"[{name}]: {resp[:600]}" for name, resp in valid[:3])
    synthesis_prompt = (
        f"Оригинальный вопрос: {prompt}\n\n"
        f"Ответы от разных AI:\n{combined}\n\n"
        f"Синтезируй ЛУЧШИЙ ответ:\n"
        f"1. Возьми самые точные факты из каждого\n"
        f"2. Убери повторения\n"
        f"3. Добавь то что упустили другие\n"
        f"4. Напиши единый чёткий ответ\n"
        f"НЕ упоминай что это синтез нескольких AI."
    )
    final, _ = await ask_groq(synthesis_prompt, "Ты синтезатор ответов. Будь точным и кратким.")
    if final:
        agent_names = "+".join(n for n, _ in valid[:3])
        return final, f"🤖 Multi-Agent ({agent_names})"

    # Если синтез не удался — берём самый длинный ответ
    best = max(valid, key=lambda x: len(x[1]))
    return best[1], best[0]


def is_complex_question(text: str) -> bool:
    """Определяет нужен ли мультиагентный подход"""
    complex_markers = [
        "сравни", "проанализируй", "объясни почему", "что лучше",
        "плюсы и минусы", "за и против", "мнения", "дискуссия",
        "спорный", "противоречие", "разные точки зрения",
        "мультиагент", "несколько ai", "все ai"
    ]
    text_low = text.lower()
    return any(m in text_low for m in complex_markers) or len(text) > 200


# ============================================================
# ТРЕКЕР ПРИВЫЧЕК
# ============================================================

async def habits_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Трекер привычек"""
    data = load_data()
    habits = data.get("habits_tracker", [])
    today = datetime.now().strftime("%d.%m.%Y")

    if not habits:
        text = (
            "🏃 *Трекер привычек*\n\n"
            "Добавь привычки и отмечай каждый день!\n\n"
            "Примеры:\n"
            "• Спорт 30 минут\n"
            "• Читать книгу\n"
            "• Пить 2л воды\n"
            "• Медитация\n\n"
            "Напиши название привычки:"
        )
        context.user_data["waiting_for"] = "habit_add"
        await update.message.reply_text(text, parse_mode="Markdown")
        return

    lines = ["🏃 *Трекер привычек*\n", f"📅 Сегодня: *{today}*\n"]
    for h in habits:
        checks = h.get("checks", {})
        done_today = checks.get(today, False)
        streak = _calc_streak(checks)
        total = sum(1 for v in checks.values() if v)
        icon = "✅" if done_today else "⬜"
        lines.append(f"{icon} *{h['name']}*")
        lines.append(f"   🔥 Серия: {streak} дн | 📊 Всего: {total} раз\n")

    text = "\n".join(lines)
    btns = []
    for i, h in enumerate(habits):
        checks = h.get("checks", {})
        done = checks.get(today, False)
        label = f"{'✅' if done else '⬜'} {h['name'][:20]}"
        btns.append([InlineKeyboardButton(label, callback_data=f"habit_check_{i}")])
    btns.append([
        InlineKeyboardButton("➕ Добавить", callback_data="habit_add"),
        InlineKeyboardButton("📊 Статистика", callback_data="habit_stats"),
    ])
    btns.append([InlineKeyboardButton("🗑 Удалить привычку", callback_data="habit_delete")])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=InlineKeyboardMarkup(btns))


def _calc_streak(checks: dict) -> int:
    """Считает текущую серию дней подряд"""
    if not checks:
        return 0
    streak = 0
    day = datetime.now()
    for _ in range(365):
        key = day.strftime("%d.%m.%Y")
        if checks.get(key):
            streak += 1
            day -= timedelta(days=1)
        else:
            break
    return streak


async def send_habit_stats(app, chat_id: int, data: dict) -> None:
    """Отправляет статистику привычек с графиком"""
    habits = data.get("habits_tracker", [])
    if not habits:
        await app.bot.send_message(chat_id, "Нет привычек для анализа.")
        return

    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches
        import numpy as np

        # Последние 14 дней
        days = []
        day = datetime.now()
        for _ in range(14):
            days.append(day.strftime("%d.%m.%Y"))
            day -= timedelta(days=1)
        days.reverse()
        day_labels = [d[:5] for d in days]  # dd.mm

        fig, axes = plt.subplots(len(habits), 1,
                                  figsize=(12, max(3, len(habits) * 1.8)),
                                  facecolor='#0f0f1a')
        if len(habits) == 1:
            axes = [axes]

        colors_done = '#00d4aa'
        colors_miss = '#1f1f35'

        for ax, h in zip(axes, habits):
            checks = h.get("checks", {})
            values = [1 if checks.get(d) else 0 for d in days]
            colors = [colors_done if v else colors_miss for v in values]
            bars = ax.bar(range(len(days)), [1]*len(days), color=colors,
                         width=0.7, edgecolor='#0f0f1a', linewidth=0.5)
            ax.set_facecolor('#0f0f1a')
            ax.set_yticks([])
            ax.set_xticks(range(len(days)))
            ax.set_xticklabels(day_labels, fontsize=7, color='#a0a0c0', rotation=45)
            streak = _calc_streak(checks)
            total = sum(values)
            pct = int(total / len(days) * 100)
            ax.set_title(f"{h['name']}   🔥{streak}д  ✅{total}/14  {pct}%",
                        color='#e8e8f5', fontsize=10, loc='left', pad=4)
            for spine in ax.spines.values():
                spine.set_visible(False)

        plt.tight_layout(pad=1.5)
        fname = f"/tmp/habits_{chat_id}.png"
        plt.savefig(fname, dpi=130, bbox_inches='tight', facecolor='#0f0f1a')
        plt.close()

        with open(fname, 'rb') as f:
            await app.bot.send_photo(
                chat_id=chat_id, photo=f,
                caption="📊 *Трекер привычек — последние 14 дней*\n🟩 выполнено  ⬛ пропущено",
                parse_mode="Markdown"
            )
        import os; os.remove(fname)

    except ImportError:
        # Текстовый fallback если нет matplotlib
        lines = ["📊 *Статистика привычек (14 дней)*\n"]
        days_list = []
        d = datetime.now()
        for _ in range(14):
            days_list.append(d.strftime("%d.%m.%Y"))
            d -= timedelta(days=1)
        days_list.reverse()

        for h in habits:
            checks = h.get("checks", {})
            bar = "".join("✅" if checks.get(day) else "⬜" for day in days_list)
            streak = _calc_streak(checks)
            total = sum(1 for day in days_list if checks.get(day))
            lines.append(f"*{h['name']}*\n{bar}\n🔥 Серия: {streak} | Всего: {total}/14\n")
        await app.bot.send_message(chat_id, "\n".join(lines), parse_mode="Markdown")


# ============================================================
# ГРАФИКИ ИЗ EXCEL / CSV
# ============================================================

async def create_chart_from_data(file_bytes: bytes, filename: str, chart_type: str = "auto") -> tuple:
    """Создаёт красивый график из Excel/CSV данных"""
    import io
    fname_lower = filename.lower()
    rows_data = []
    headers = []

    try:
        if fname_lower.endswith(".csv"):
            import csv
            text = file_bytes.decode("utf-8", errors="ignore")
            reader = csv.reader(io.StringIO(text))
            for i, row in enumerate(reader):
                if i == 0:
                    headers = row
                elif i <= 100:
                    rows_data.append(row)
        elif fname_lower.endswith((".xlsx", ".xls")):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            ws = wb.active
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i == 0:
                    headers = [str(c) if c else f"Col{j}" for j, c in enumerate(row)]
                elif i <= 100:
                    rows_data.append(list(row))
    except Exception as e:
        return None, f"❌ Ошибка чтения: {e}"

    if not rows_data or not headers:
        return None, "❌ Файл пустой или не удалось прочитать данные."

    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        # Находим числовые колонки
        numeric_cols = []
        label_col = None
        for j, col in enumerate(headers):
            values = []
            for row in rows_data:
                try:
                    v = float(str(row[j]).replace(",", ".").replace(" ", "")) if j < len(row) else None
                    if v is not None:
                        values.append(v)
                except:
                    pass
            if len(values) > len(rows_data) * 0.5:
                numeric_cols.append((j, col, values))
            elif label_col is None and j < 3:
                label_col = j

        if not numeric_cols:
            return None, "❌ Не найдено числовых данных для графика."

        labels = [str(row[label_col]) if label_col is not None and label_col < len(row) else str(i+1)
                  for i, row in enumerate(rows_data)][:30]

        fig, ax = plt.subplots(figsize=(12, 6), facecolor='#0f0f1a')
        ax.set_facecolor('#0f0f1a')

        palette = ['#6c5ce7', '#a855f7', '#00d4aa', '#fd79a8', '#fdcb6e', '#74b9ff']

        if len(numeric_cols) == 1:
            # Один ряд — столбчатый или линейный
            j, col, values = numeric_cols[0]
            vals = values[:len(labels)]
            colors = [palette[i % len(palette)] for i in range(len(vals))]
            ax.bar(range(len(vals)), vals, color=colors, edgecolor='#0f0f1a', linewidth=0.5)
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels, rotation=45, ha='right', fontsize=8, color='#a0a0c0')
            ax.set_title(col, color='#e8e8f5', fontsize=13, pad=12)
        else:
            # Несколько рядов — линейный
            x = range(len(labels))
            for idx, (j, col, values) in enumerate(numeric_cols[:5]):
                vals = values[:len(labels)]
                color = palette[idx % len(palette)]
                ax.plot(x, vals, color=color, linewidth=2, marker='o', markersize=4, label=col)
            ax.set_xticks(list(x))
            ax.set_xticklabels(labels, rotation=45, ha='right', fontsize=8, color='#a0a0c0')
            ax.legend(facecolor='#1a1a2e', labelcolor='#e8e8f5', fontsize=9)

        ax.tick_params(colors='#a0a0c0')
        ax.yaxis.label.set_color('#a0a0c0')
        for spine in ax.spines.values():
            spine.set_color('#2a2a4a')
        ax.grid(axis='y', color='#2a2a4a', linestyle='--', alpha=0.5)

        plt.tight_layout()
        fname_out = f"/tmp/chart_{id(file_bytes)}.png"
        plt.savefig(fname_out, dpi=140, bbox_inches='tight', facecolor='#0f0f1a')
        plt.close()
        return fname_out, None

    except ImportError:
        return None, "❌ Библиотека matplotlib не установлена."
    except Exception as e:
        return None, f"❌ Ошибка создания графика: {e}"


# ============================================================
# ФИЧА 1: ДЕБАТЫ — бот спорит с тобой
# ============================================================

DEBATE_SESSIONS = {}  # chat_id -> {topic, position, round}

async def debate_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Режим дебатов"""
    text = (
        "🤝 *Режим дебатов*\n\n"
        "Я займу противоположную позицию и буду спорить с тобой!\n\n"
        "Это прокачивает:\n"
        "• Критическое мышление\n"
        "• Умение аргументировать\n"
        "• Видеть разные стороны вопроса\n\n"
        "Напиши тему для дебатов.\nПример: _AI заменит всех программистов_"
    )
    context.user_data["waiting_for"] = "debate_topic"
    await update.message.reply_text(text, parse_mode="Markdown")


async def run_debate(app, chat_id: int, topic: str, user_arg: str, round_num: int) -> None:
    """Проводит раунд дебатов"""
    try:
        is_first = round_num == 1
        system = (
            "Ты участник дебатов. Занимаешь ПРОТИВОПОЛОЖНУЮ позицию пользователю. "
            "Твоя задача — убедительно спорить, приводить факты, цифры, примеры. "
            "Будь напористым но уважительным. Заканчивай аргумент вопросом к оппоненту."
        )
        if is_first:
            prompt = (
                f"Тема дебатов: {topic}\n\n"
                f"Пользователь занял позицию ЗА. Ты занимаешь позицию ПРОТИВ.\n"
                f"Его аргумент: {user_arg}\n\n"
                f"Раунд 1: Дай 2-3 сильных контраргумента с фактами. "
                f"В конце задай острый вопрос оппоненту."
            )
        else:
            prompt = (
                f"Тема: {topic} | Раунд {round_num}\n"
                f"Аргумент оппонента: {user_arg}\n\n"
                f"Разбей его аргумент, приведи новые факты, усиль свою позицию. "
                f"В конце подведи мини-итог раунда и задай следующий вопрос."
            )
        response, ai_name = await get_best_ai_response(prompt, system)
        if response:
            header = f"⚔️ *Раунд {round_num} — Контраргумент:*\n\n" if not is_first else f"⚔️ *Дебаты: {topic}*\n*Мой аргумент ПРОТИВ:*\n\n"
            await send_long_message(app.bot, chat_id, header + response, parse_mode="Markdown")
            if round_num >= 5:
                # Финал — подводим итог
                final_prompt = (
                    f"Дебаты по теме '{topic}' завершены после 5 раундов.\n"
                    f"Подведи объективный итог: кто привёл более сильные аргументы, "
                    f"что было убедительно с обеих сторон, и какой общий вывод."
                )
                final, _ = await get_best_ai_response(final_prompt, "Ты объективный судья дебатов.")
                if final:
                    await send_long_message(app.bot, chat_id, f"🏆 *Итог дебатов:*\n\n{final}", parse_mode="Markdown")
                await app.bot.send_message(chat_id, "Дебаты завершены! Напиши /debate чтобы начать новые.")
                DEBATE_SESSIONS.pop(chat_id, None)
    except Exception as e:
        logger.error(f"debate error: {e}")


# ============================================================
# ФИЧА 2: ЭКСПОРТ ВСЕГО В PDF
# ============================================================

async def export_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Экспорт всех данных в PDF"""
    msg = await update.message.reply_text("📄 Создаю PDF со всеми твоими данными...")
    data = load_data()

    try:
        lines = []
        now = datetime.now().strftime("%d.%m.%Y %H:%M")
        profile = data.get("profile", {})
        name = profile.get("name", "Пользователь")

        lines.append(f"ЛИЧНЫЙ ДАЙДЖЕСТ — {name}")
        lines.append(f"Дата экспорта: {now}")
        lines.append("=" * 60)

        # Профиль
        lines.append("\n[ПРОФИЛЬ]")
        lines.append(f"Имя: {name}")
        lines.append(f"Город: {profile.get('city', 'не указан')}")
        if profile.get("interests"):
            lines.append(f"Интересы: {', '.join(profile['interests'])}")
        if profile.get("goals"):
            lines.append(f"Цели: {', '.join(profile['goals'])}")

        # Статистика
        s = data.get("stats", {})
        lines.append("\n[СТАТИСТИКА]")
        lines.append(f"Сообщений всего: {s.get('messages_total', 0)}")
        lines.append(f"Голосовых: {s.get('voice_total', 0)}")
        lines.append(f"Файлов: {s.get('files_total', 0)}")
        lines.append(f"В боте с: {s.get('first_seen', 'неизвестно')}")

        # Цели
        goals = data.get("goals_system", {}).get("goals", [])
        if goals:
            lines.append("\n[ЦЕЛИ]")
            for g in goals:
                status = "✓" if g.get("status") == "done" else "○"
                lines.append(f"{status} {g['title']} — {g.get('progress', 0)}% [{g.get('deadline', 'без срока')}]")

        # Задачи
        tasks = data.get("tasks", [])
        if tasks:
            active = [t for t in tasks if not t.get("done")]
            done = [t for t in tasks if t.get("done")]
            lines.append("\n[ЗАДАЧИ]")
            lines.append(f"Активных: {len(active)} | Выполнено: {len(done)}")
            for t in active[:20]:
                lines.append(f"  ○ {t['text']}")
            for t in done[:10]:
                lines.append(f"  ✓ {t['text']}")

        # Привычки
        habits = data.get("habits_tracker", [])
        if habits:
            lines.append("\n[ПРИВЫЧКИ]")
            for h in habits:
                checks = h.get("checks", {})
                streak = _calc_streak(checks)
                total = sum(1 for v in checks.values() if v)
                lines.append(f"• {h['name']} — серия: {streak} дн, всего: {total} раз")

        # База знаний
        brain = data.get("second_brain", {})
        knowledge = brain.get("knowledge", {})
        facts = brain.get("facts", [])
        if knowledge:
            lines.append("\n[БАЗА ЗНАНИЙ]")
            for k, v in list(knowledge.items())[:30]:
                lines.append(f"• {v['title']} ({len(v.get('notes', []))} заметок)")
        if facts:
            lines.append("\n[ФАКТЫ ИЗ ПАМЯТИ]")
            for f in facts[:30]:
                lines.append(f"• {f}")

        # Изученные темы
        topics = data.get("topics", [])
        if topics:
            lines.append("\n[ИЗУЧЕННЫЕ ТЕМЫ]")
            for t in topics[-50:]:
                lines.append(f"• {t}")

        # История последнего чата
        history = data.get("chats", {}).get(data.get("current", "main"), {}).get("history", [])
        if history:
            lines.append("\n[ПОСЛЕДНИЕ СООБЩЕНИЯ]")
            for msg_h in history[-20:]:
                role = msg_h.get("role", "?")
                text = msg_h.get("text", "")[:200]
                lines.append(f"{role}: {text}")

        # Создаём PDF через reportlab или txt fallback
        content_text = "\n".join(lines)
        ts = datetime.now().strftime("%d%m%Y_%H%M")
        fname = f"/tmp/export_{ts}.txt"

        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.lib.units import cm
            import textwrap

            pdf_fname = f"/tmp/export_{ts}.pdf"
            c = canvas.Canvas(pdf_fname, pagesize=A4)
            w, h = A4
            y = h - 2*cm
            c.setFont("Helvetica-Bold", 16)
            c.drawString(2*cm, y, f"Личный дайджест — {name}")
            y -= 0.8*cm
            c.setFont("Helvetica", 9)
            c.drawString(2*cm, y, f"Экспортировано: {now}")
            y -= 0.5*cm

            for line in lines[3:]:
                if y < 2*cm:
                    c.showPage()
                    y = h - 2*cm
                if line.startswith("[") and line.endswith("]"):
                    c.setFont("Helvetica-Bold", 11)
                    y -= 0.3*cm
                elif line.startswith("="):
                    c.setFont("Helvetica", 8)
                    c.drawString(2*cm, y, "-" * 80)
                    y -= 0.4*cm
                    continue
                else:
                    c.setFont("Helvetica", 9)

                wrapped = textwrap.wrap(line, 90) or [""]
                for wl in wrapped:
                    c.drawString(2*cm, y, wl)
                    y -= 0.4*cm
            c.save()
            export_file = pdf_fname
            export_name = f"Дайджест_{name}_{ts}.pdf"
        except ImportError:
            with open(fname, "w", encoding="utf-8") as f:
                f.write(content_text)
            export_file = fname
            export_name = f"Дайджест_{name}_{ts}.txt"

        await msg.delete()
        with open(export_file, "rb") as f:
            await context.bot.send_document(
                chat_id=update.effective_chat.id,
                document=f,
                filename=export_name,
                caption=f"📄 *Твой личный дайджест*\n\n"
                        f"📊 {s.get('messages_total',0)} сообщений\n"
                        f"🎯 {len(goals)} целей\n"
                        f"✅ {len(tasks)} задач\n"
                        f"📚 {len(knowledge)} тем",
                parse_mode="Markdown"
            )
        import os; os.remove(export_file)

    except Exception as e:
        logger.error(f"export error: {e}")
        await msg.edit_text(f"❌ Ошибка экспорта: {e}")


# ============================================================
# ФИЧА 3: СИМУЛЯТОР СОБЕСЕДОВАНИЙ
# ============================================================

INTERVIEW_SESSIONS = {}

async def interview_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Симулятор собеседований"""
    text = (
        "💼 *Симулятор собеседования*\n\n"
        "Я сыграю роль интервьюера!\n\n"
        "Доступные режимы:\n"
        "• _Python разработчик_\n"
        "• _Frontend разработчик_\n"
        "• _Product Manager_\n"
        "• _Data Scientist_\n"
        "• _Маркетолог_\n"
        "• Или напиши свою профессию\n\n"
        "На какую позицию собеседуемся?"
    )
    context.user_data["waiting_for"] = "interview_role"
    await update.message.reply_text(text, parse_mode="Markdown")


async def run_interview(app, chat_id: int, role: str, answer: str, q_num: int) -> None:
    """Проводит раунд собеседования"""
    try:
        if q_num == 1:
            prompt = (
                f"Ты строгий интервьюер в крупной IT компании. Проводишь собеседование на позицию: {role}\n\n"
                f"Начни собеседование:\n"
                f"1. Представься как интервьюер\n"
                f"2. Задай первый вопрос (начни с простого — расскажи о себе)\n"
                f"Будь профессиональным и немного строгим."
            )
            system = f"Ты опытный HR и технический интервьюер для позиции {role}."
        else:
            prompt = (
                f"Позиция: {role} | Вопрос №{q_num}\n\n"
                f"Ответ кандидата: {answer}\n\n"
                f"1. Дай краткую оценку ответа (что хорошо, что можно улучшить)\n"
                f"2. Задай следующий вопрос (более сложный)\n"
                f"После 8 вопросов — подведи итог собеседования с оценкой по 10-балльной шкале."
            )
            system = f"Ты опытный интервьюер. Оценивай честно и конструктивно. Позиция: {role}"

        response, _ = await get_best_ai_response(prompt, system)
        if response:
            prefix = f"💼 *Вопрос {q_num}:*\n\n" if q_num > 1 else "💼 *Собеседование началось!*\n\n"
            await send_long_message(app.bot, chat_id, prefix + response, parse_mode="Markdown")
            if q_num >= 8:
                INTERVIEW_SESSIONS.pop(chat_id, None)
                await app.bot.send_message(chat_id, "✅ Собеседование завершено!\nНапиши /interview чтобы попробовать снова.")
    except Exception as e:
        logger.error(f"interview error: {e}")


# ============================================================
# ФИЧА 4: КНИЖНЫЙ КЛУБ
# ============================================================

async def books_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Книжный клуб"""
    data = load_data()
    books = data.get("books", [])

    if not books:
        text = (
            "📖 *Книжный клуб*\n\n"
            "Добавь книгу и я сделаю:\n"
            "• 📌 Краткое содержание\n"
            "• 🔑 10 ключевых идей\n"
            "• 💬 Вопросы для размышления\n"
            "• ⭐ Главные цитаты\n"
            "• 📚 Похожие книги\n\n"
            "Напиши название книги и автора:"
        )
        context.user_data["waiting_for"] = "book_add"
    else:
        lines = ["📖 *Книжный клуб*\n"]
        reading = [b for b in books if b.get("status") == "reading"]
        done = [b for b in books if b.get("status") == "done"]
        want = [b for b in books if b.get("status") == "want"]
        if reading:
            lines.append("📖 *Читаю сейчас:*")
            for b in reading:
                lines.append(f"  • {b['title']} — {b.get('author','')}")
        if done:
            lines.append(f"\n✅ *Прочитано ({len(done)}):*")
            for b in done[-5:]:
                stars = "⭐" * b.get("rating", 0)
                lines.append(f"  • {b['title']} {stars}")
        if want:
            lines.append(f"\n📋 *Хочу прочитать ({len(want)}):*")
            for b in want[:5]:
                lines.append(f"  • {b['title']}")
        text = "\n".join(lines)

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("➕ Добавить книгу", callback_data="book_add"),
         InlineKeyboardButton("📊 Статистика", callback_data="book_stats")],
        [InlineKeyboardButton("🔍 Анализ книги", callback_data="book_analyze")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


async def analyze_book(app, chat_id: int, book_title: str) -> None:
    """Анализирует книгу через AI"""
    try:
        await app.bot.send_message(chat_id, f"📖 Анализирую: *{book_title}*...", parse_mode="Markdown")
        prompt = (
            f"Книга: {book_title}\n\n"
            f"Сделай полный анализ:\n"
            f"## 📌 О чём книга (3-4 предложения)\n"
            f"## 🔑 10 ключевых идей (кратко каждая)\n"
            f"## 💬 5 вопросов для размышления\n"
            f"## ⭐ 3 главные цитаты или мысли автора\n"
            f"## 👥 Кому рекомендую и почему\n"
            f"## 📚 3 похожие книги"
        )
        result, ai_name = await get_best_ai_response(
            prompt, "Ты литературный критик и книжный эксперт. Давай глубокий анализ."
        )
        if result:
            await send_long_message(
                app.bot, chat_id,
                f"📖 *Анализ: {book_title}*\n_{ai_name}_\n\n{result}",
                parse_mode="Markdown"
            )
    except Exception as e:
        logger.error(f"analyze_book error: {e}")


# ============================================================
# ФИЧА 5: МОНИТОРИНГ САЙТОВ
# ============================================================

MONITOR_FILE = "monitors.json"

def load_monitors() -> list:
    try:
        if Path(MONITOR_FILE).exists():
            with open(MONITOR_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except: pass
    return []

def save_monitors(monitors: list) -> None:
    try:
        with open(MONITOR_FILE, "w", encoding="utf-8") as f:
            json.dump(monitors, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"save_monitors error: {e}")


async def monitor_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Мониторинг сайтов"""
    monitors = load_monitors()
    active = [m for m in monitors if m.get("chat_id") == update.effective_chat.id]

    if not active:
        text = (
            "🌐 *Мониторинг сайтов*\n\n"
            "Я буду следить за сайтом и уведомлю когда что-то изменится!\n\n"
            "Примеры:\n"
            "• Цена товара на сайте\n"
            "• Новости на странице\n"
            "• Любые изменения текста\n\n"
            "Отправь ссылку для мониторинга:"
        )
        context.user_data["waiting_for"] = "monitor_add"
    else:
        lines = ["🌐 *Мониторинг сайтов*\n", f"Отслеживаю {len(active)} сайтов:\n"]
        for m in active:
            status = "✅" if m.get("active") else "⏸"
            lines.append(f"{status} {m['url'][:50]}")
            if m.get("last_check"):
                lines.append(f"   Последняя проверка: {m['last_check']}")
        text = "\n".join(lines)

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("➕ Добавить сайт", callback_data="monitor_add"),
         InlineKeyboardButton("🗑 Удалить", callback_data="monitor_delete")],
        [InlineKeyboardButton("🔄 Проверить сейчас", callback_data="monitor_check")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


async def check_monitors_loop(app) -> None:
    """Фоновый цикл проверки изменений на сайтах каждые 30 минут"""
    while True:
        try:
            await asyncio.sleep(1800)  # 30 минут
            monitors = load_monitors()
            changed = False
            for m in monitors:
                if not m.get("active"):
                    continue
                try:
                    new_content = await fetch_url_content(m["url"])
                    if new_content.startswith("error:"):
                        continue
                    old_hash = m.get("content_hash", "")
                    import hashlib
                    new_hash = hashlib.md5(new_content[:5000].encode()).hexdigest()
                    if old_hash and new_hash != old_hash:
                        m["content_hash"] = new_hash
                        m["last_check"] = datetime.now().strftime("%d.%m.%Y %H:%M")
                        changed = True
                        await app.bot.send_message(
                            m["chat_id"],
                            f"🔔 *Изменение на сайте!*\n\n"
                            f"🌐 {m['url'][:60]}\n\n"
                            f"Открой ссылку чтобы посмотреть что изменилось.",
                            parse_mode="Markdown"
                        )
                    elif not old_hash:
                        m["content_hash"] = new_hash
                        m["last_check"] = datetime.now().strftime("%d.%m.%Y %H:%M")
                        changed = True
                except Exception as e:
                    logger.error(f"monitor check error: {e}")
            if changed:
                save_monitors(monitors)
        except Exception as e:
            logger.error(f"monitor loop error: {e}")
            await asyncio.sleep(300)


# ============================================================
# АГЕНТНЫЙ РЕЖИМ — автономное выполнение сложных задач
# ============================================================

AGENT_SESSIONS = {}  # chat_id -> {task, steps, current_step, results}

async def agent_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Агентный режим — бот сам планирует и выполняет задачу"""
    text = (
        "🤖 *Агентный режим*\n\n"
        "Опиши задачу — я сам:\n"
        "1️⃣ Разобью на шаги\n"
        "2️⃣ Выполню каждый шаг\n"
        "3️⃣ Использую поиск, анализ, код\n"
        "4️⃣ Дам итоговый результат\n\n"
        "Примеры:\n"
        "• _Исследуй рынок криптовалют и дай топ-3 для инвестиций_\n"
        "• _Напиши бизнес-план для онлайн школы_\n"
        "• _Проанализируй конкурентов Telegram ботов и найди нишу_\n\n"
        "Напиши задачу:"
    )
    context.user_data["waiting_for"] = "agent_task"
    await update.message.reply_text(text, parse_mode="Markdown")


async def run_agent(app, chat_id: int, task: str, data: dict) -> None:
    """Автономно выполняет задачу в несколько шагов"""
    try:
        await app.bot.send_message(
            chat_id,
            f"🤖 *Агент запущен*\n\n📋 Задача: _{task}_\n\n⚙️ Планирую шаги...",
            parse_mode="Markdown"
        )

        # Шаг 1: Планирование
        plan_prompt = (
            f"Задача: {task}\n\n"
            f"Создай план из 4-6 конкретных шагов для выполнения этой задачи.\n"
            f"Для каждого шага укажи тип: [ПОИСК], [АНАЛИЗ], [КОД], [ТЕКСТ], [РАСЧЁТ]\n\n"
            f"Формат:\n"
            f"Шаг 1 [ТИП]: описание\n"
            f"Шаг 2 [ТИП]: описание\n"
            f"...\n\n"
            f"Только план, без лишних слов."
        )
        plan_raw, _ = await get_best_ai_response(
            plan_prompt,
            "Ты агент-планировщик. Создавай чёткие выполнимые планы."
        )

        if not plan_raw:
            await app.bot.send_message(chat_id, "❌ Не удалось создать план.")
            return

        # Парсим шаги
        import re
        steps = re.findall(r'Шаг \d+\s*\[([^\]]+)\]:\s*(.+)', plan_raw)
        if not steps:
            steps = [(f"АНАЛИЗ", line.strip()) for line in plan_raw.split('\n') if line.strip() and len(line) > 10]
        steps = steps[:6]

        plan_text = "\n".join(f"{'🔍' if 'ПОИСК' in s[0] else '🧠' if 'АНАЛИЗ' in s[0] else '💻' if 'КОД' in s[0] else '📝'} {s[1]}" for s in steps)
        await app.bot.send_message(
            chat_id,
            f"📋 *План ({len(steps)} шагов):*\n\n{plan_text}\n\n▶️ Начинаю выполнение...",
            parse_mode="Markdown"
        )

        # Шаг 2: Выполнение каждого шага
        results = []
        context_so_far = f"Задача: {task}\n\n"

        for i, (step_type, step_desc) in enumerate(steps, 1):
            step_msg = await app.bot.send_message(
                chat_id,
                f"⚙️ *Шаг {i}/{len(steps)}:* _{step_desc}_",
                parse_mode="Markdown"
            )

            step_result = ""

            try:
                if "ПОИСК" in step_type.upper():
                    # Поиск в интернете
                    search_result = await search_web(step_desc)
                    step_prompt = (
                        f"Контекст задачи: {task}\n"
                        f"Текущий шаг: {step_desc}\n"
                        f"Данные из поиска: {search_result[:3000]}\n\n"
                        f"Проанализируй эти данные применительно к задаче. Выдели ключевые факты."
                    )
                    step_result, _ = await get_best_ai_response(step_prompt, "Ты аналитик. Будь кратким и конкретным.")

                elif "КОД" in step_type.upper():
                    code_prompt = (
                        f"Задача: {task}\n"
                        f"Нужно: {step_desc}\n"
                        f"Контекст: {context_so_far[-1000:]}\n\n"
                        f"Напиши код для решения этого шага."
                    )
                    step_result, _ = await get_best_ai_response(code_prompt, "Ты разработчик. Пиши рабочий код.")

                else:
                    # Анализ / Текст / Расчёт
                    analysis_prompt = (
                        f"Задача: {task}\n"
                        f"Выполни шаг: {step_desc}\n\n"
                        f"Уже известно:\n{context_so_far[-2000:]}\n\n"
                        f"Дай конкретный результат для этого шага."
                    )
                    step_result, _ = await get_best_ai_response(
                        analysis_prompt,
                        "Ты умный агент. Выполняй шаги конкретно и детально."
                    )

            except Exception as e:
                step_result = f"Ошибка на шаге: {e}"

            results.append(f"[Шаг {i}: {step_desc}]\n{step_result or 'нет данных'}")
            context_so_far += f"\nШаг {i} ({step_desc}):\n{(step_result or '')[:500]}\n"

            # Обновляем сообщение о шаге
            short = (step_result or "")[:150] + ("..." if len(step_result or "") > 150 else "")
            try:
                await app.bot.edit_message_text(
                    chat_id=chat_id,
                    message_id=step_msg.message_id,
                    text=f"✅ *Шаг {i}/{len(steps)}:* _{step_desc}_\n\n{short}",
                    parse_mode="Markdown"
                )
            except: pass
            await asyncio.sleep(0.5)

        # Шаг 3: Итоговый синтез
        await app.bot.send_message(chat_id, "🧠 Синтезирую итоговый результат...")

        all_results = "\n\n".join(results)
        final_prompt = (
            f"Задача: {task}\n\n"
            f"Результаты всех шагов:\n{all_results[:6000]}\n\n"
            f"Синтезируй ИТОГОВЫЙ ОТЧЁТ:\n"
            f"## 🎯 ОТВЕТ НА ЗАДАЧУ\n"
            f"## 📊 КЛЮЧЕВЫЕ ВЫВОДЫ\n"
            f"## ✅ РЕКОМЕНДАЦИИ\n"
            f"## 🚀 СЛЕДУЮЩИЕ ШАГИ\n\n"
            f"Будь конкретным, используй данные из шагов."
        )
        final_result, ai_name = await get_best_ai_response(
            final_prompt,
            "Ты опытный аналитик. Синтезируй чёткий итог из всех данных."
        )

        if final_result:
            # Сохраняем в базу знаний
            await save_to_brain(data, f"Агент: {task[:40]}", final_result[:500])
            save_data(data)

            # Отправляем как файл если длинный
            if len(final_result) > 2000:
                ts = datetime.now().strftime("%d%m%Y_%H%M")
                fname = f"/tmp/agent_{ts}.txt"
                full_content = (
                    f"АГЕНТНЫЙ ОТЧЁТ\n{'='*50}\n"
                    f"Задача: {task}\n"
                    f"Дата: {datetime.now().strftime('%d.%m.%Y %H:%M')}\n"
                    f"{'='*50}\n\n"
                    f"ПЛАН ВЫПОЛНЕНИЯ:\n{plan_text}\n\n"
                    f"{'='*50}\nИТОГОВЫЙ ОТЧЁТ:\n{'='*50}\n\n"
                    f"{final_result}\n\n"
                    f"{'='*50}\nДЕТАЛЬНЫЕ РЕЗУЛЬТАТЫ ШАГОВ:\n{'='*50}\n\n"
                    f"{all_results}"
                )
                with open(fname, "w", encoding="utf-8") as f:
                    f.write(full_content)
                with open(fname, "rb") as f:
                    await app.bot.send_document(
                        chat_id=chat_id, document=f,
                        filename=f"Агент_{task[:25]}_{ts}.txt",
                        caption=f"🤖 *Агент завершил задачу!*\n_{ai_name}_\n\n📊 Шагов выполнено: {len(steps)}",
                        parse_mode="Markdown"
                    )
                import os; os.remove(fname)
            else:
                await send_long_message(
                    app.bot, chat_id,
                    f"🤖 *Агент завершил задачу!*\n_{ai_name}_\n\n{final_result}",
                    parse_mode="Markdown"
                )
        else:
            await app.bot.send_message(chat_id, "❌ Не удалось синтезировать итог.")

    except Exception as e:
        logger.error(f"agent error: {e}")
        await app.bot.send_message(chat_id, f"❌ Ошибка агента: {e}")


# ============================================================
# ФИЧА 1: EMAIL ИНТЕГРАЦИЯ (Gmail через IMAP/SMTP)
# ============================================================

async def email_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Email интеграция"""
    data = load_data()
    email_cfg = data.get("email_config", {})

    if not email_cfg.get("address"):
        text = (
            "📧 *Email интеграция*\n\n"
            "Подключи Gmail и я смогу:\n"
            "• 📬 Читать новые письма\n"
            "• ✉️ Отправлять письма голосом\n"
            "• 📋 Делать саммари писем\n"
            "• 🔔 Уведомлять о важных письмах\n\n"
            "Для подключения нужен Gmail + App Password.\n\n"
            "📌 Как получить App Password:\n"
            "1. Зайди в Google Account\n"
            "2. Безопасность → Двухэтапная аутентификация\n"
            "3. Пароли приложений → Создать\n\n"
            "Введи свой Gmail адрес:"
        )
        context.user_data["waiting_for"] = "email_address"
    else:
        addr = email_cfg.get("address", "")
        text = f"📧 *Email подключён*\n\n📬 {addr}\n\nЧто сделать?"

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("📬 Новые письма", callback_data="email_inbox"),
         InlineKeyboardButton("✉️ Написать письмо", callback_data="email_compose")],
        [InlineKeyboardButton("📋 Саммари входящих", callback_data="email_summary"),
         InlineKeyboardButton("⚙️ Настройки", callback_data="email_settings")],
    ])
    await update.message.reply_text(text, parse_mode="Markdown",
                                     reply_markup=keyboard if email_cfg.get("address") else None)


async def fetch_emails(email_addr: str, app_password: str, count: int = 5) -> list:
    """Читает последние письма через IMAP"""
    try:
        import imaplib, email as email_lib
        from email.header import decode_header

        mail = imaplib.IMAP4_SSL("imap.gmail.com")
        mail.login(email_addr, app_password)
        mail.select("inbox")

        _, data = mail.search(None, "UNSEEN")
        ids = data[0].split()[-count:]

        emails = []
        for eid in reversed(ids):
            _, msg_data = mail.fetch(eid, "(RFC822)")
            msg = email_lib.message_from_bytes(msg_data[0][1])

            subject_raw, enc = decode_header(msg["Subject"] or "")[0]
            subject = subject_raw.decode(enc or "utf-8") if isinstance(subject_raw, bytes) else subject_raw

            from_raw, enc2 = decode_header(msg["From"] or "")[0]
            sender = from_raw.decode(enc2 or "utf-8") if isinstance(from_raw, bytes) else from_raw

            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_payload(decode=True).decode("utf-8", errors="ignore")[:500]
                        break
            else:
                body = msg.get_payload(decode=True).decode("utf-8", errors="ignore")[:500]

            emails.append({
                "subject": subject, "from": sender,
                "date": msg["Date"], "body": body
            })
        mail.logout()
        return emails
    except Exception as e:
        return [{"error": str(e)}]


async def send_email(email_addr: str, app_password: str, to: str, subject: str, body: str) -> bool:
    """Отправляет письмо через SMTP"""
    try:
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart

        msg = MIMEMultipart()
        msg["From"] = email_addr
        msg["To"] = to
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain", "utf-8"))

        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
            server.login(email_addr, app_password)
            server.send_message(msg)
        return True
    except Exception as e:
        logger.error(f"send_email error: {e}")
        return False


# ============================================================
# ФИЧА 2: МАРШРУТЫ В КАРТАХ
# ============================================================

async def get_route_info(from_place: str, to_place: str) -> str:
    """Получает информацию о маршруте через Nominatim"""
    try:
        import urllib.parse

        async def geocode(place: str) -> tuple:
            url = f"https://nominatim.openstreetmap.org/search?q={urllib.parse.quote(place)}&format=json&limit=1"
            async with httpx.AsyncClient() as client:
                r = await client.get(url, headers={"User-Agent": "AIAgentBot/1.0"}, timeout=10)
                data = r.json()
            if data:
                return float(data[0]["lat"]), float(data[0]["lon"]), data[0]["display_name"]
            return None, None, None

        lat1, lon1, name1 = await geocode(from_place)
        lat2, lon2, name2 = await geocode(to_place)

        if not lat1 or not lat2:
            return f"❌ Не удалось найти один из адресов."

        # Расстояние по прямой (Haversine)
        import math
        R = 6371
        dlat = math.radians(lat2 - lat1)
        dlon = math.radians(lon2 - lon1)
        a = math.sin(dlat/2)**2 + math.cos(math.radians(lat1)) * math.cos(math.radians(lat2)) * math.sin(dlon/2)**2
        dist_km = R * 2 * math.asin(math.sqrt(a))

        # Примерное время
        walk_min = int(dist_km / 0.083)  # 5 км/ч
        car_min = int(dist_km / 0.5)     # 30 км/ч в городе
        transit_min = int(dist_km / 0.4)  # 24 км/ч

        map_url = f"https://www.openstreetmap.org/directions?from={lat1},{lon1}&to={lat2},{lon2}"

        return (
            f"🗺 *Маршрут*\n\n"
            f"📍 Откуда: {name1[:60]}\n"
            f"🏁 Куда: {name2[:60]}\n\n"
            f"📏 Расстояние: *{dist_km:.1f} км*\n\n"
            f"⏱ Примерное время:\n"
            f"🚗 На машине: ~{car_min} мин\n"
            f"🚌 Транспорт: ~{transit_min} мин\n"
            f"🚶 Пешком: ~{walk_min} мин\n\n"
            f"🔗 [Открыть маршрут на карте]({map_url})"
        )
    except Exception as e:
        return f"❌ Ошибка построения маршрута: {e}"


def is_route_request(text: str) -> tuple:
    """Определяет запрос маршрута"""
    import re
    patterns = [
        r'маршрут (?:из|от)\s+(.+?)\s+до\s+(.+)',
        r'как (?:доехать|добраться|дойти) (?:из|от)\s+(.+?)\s+до\s+(.+)',
        r'дорога (?:из|от)\s+(.+?)\s+до\s+(.+)',
        r'расстояние (?:от|из)\s+(.+?)\s+до\s+(.+)',
    ]
    for p in patterns:
        m = re.search(p, text.lower())
        if m:
            return True, m.group(1).strip(), m.group(2).strip()
    return False, "", ""


# ============================================================
# ФИЧА 3: FINE-TUNING — персонализация стиля бота
# ============================================================

async def finetune_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Персонализация стиля общения бота"""
    data = load_data()
    style = data.get("user_style", {})

    if not style.get("analyzed"):
        text = (
            "🧠 *Персонализация стиля*\n\n"
            "Я изучу как ты общаешься и начну отвечать в твоём стиле!\n\n"
            "Анализирую:\n"
            "• Длину твоих сообщений\n"
            "• Стиль (формальный/неформальный)\n"
            "• Любимые темы\n"
            "• Как ты задаёшь вопросы\n"
            "• Предпочтения в ответах\n\n"
            "Нажми кнопку — я проанализирую историю наших разговоров!"
        )
        keyboard = InlineKeyboardMarkup([
            [InlineKeyboardButton("🔍 Анализировать мой стиль", callback_data="style_analyze")],
        ])
    else:
        tone = style.get("tone", "нейтральный")
        length = style.get("preferred_length", "средний")
        topics = style.get("top_topics", [])
        text = (
            f"🧠 *Мой стиль общения*\n\n"
            f"🎭 Тон: *{tone}*\n"
            f"📏 Длина ответов: *{length}*\n"
            f"📚 Топ темы: {', '.join(topics[:5]) if topics else 'анализируется'}\n"
            f"✅ Бот адаптирован под тебя!\n\n"
            f"_Чем больше общаешься — тем точнее персонализация_"
        )
        keyboard = InlineKeyboardMarkup([
            [InlineKeyboardButton("🔄 Обновить анализ", callback_data="style_analyze"),
             InlineKeyboardButton("🔧 Настроить вручную", callback_data="style_manual")],
            [InlineKeyboardButton("📊 Подробный отчёт", callback_data="style_report")],
        ])
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=keyboard)


async def analyze_user_style(data: dict) -> dict:
    """Анализирует стиль пользователя по истории чатов"""
    try:
        all_messages = []
        for chat_id, chat in data.get("chats", {}).items():
            for msg in chat.get("history", []):
                if msg.get("role") == "Пользователь":
                    all_messages.append(msg.get("text", ""))

        if len(all_messages) < 5:
            return {"analyzed": False, "reason": "мало сообщений"}

        # Базовый анализ без AI
        total_len = sum(len(m) for m in all_messages)
        avg_len = total_len / len(all_messages)
        short_msgs = sum(1 for m in all_messages if len(m) < 30)
        long_msgs = sum(1 for m in all_messages if len(m) > 100)

        # Тон
        informal_words = ["привет", "окей", "ок", "лол", "давай", "круто", "норм", "спс"]
        formal_words = ["пожалуйста", "благодарю", "прошу", "необходимо", "следует"]
        informal_count = sum(1 for m in all_messages for w in informal_words if w in m.lower())
        formal_count = sum(1 for m in all_messages for w in formal_words if w in m.lower())
        tone = "неформальный" if informal_count > formal_count else "формальный" if formal_count > informal_count else "нейтральный"

        # Предпочитаемая длина ответов
        if avg_len < 40:
            preferred_length = "краткий"
        elif avg_len > 120:
            preferred_length = "подробный"
        else:
            preferred_length = "средний"

        # Любимые темы из базы знаний
        topics = data.get("topics", [])[-10:]
        brain_topics = list(data.get("second_brain", {}).get("knowledge", {}).keys())[:5]
        all_topics = list(set(topics + brain_topics))

        # AI анализ стиля
        sample = "\n".join(all_messages[-20:])[:2000]
        style_prompt = (
            f"Проанализируй стиль общения пользователя по его сообщениям:\n\n{sample}\n\n"
            f"Дай краткую характеристику в 3-4 предложениях:\n"
            f"- Как он общается (формально/неформально)\n"
            f"- Какие вопросы задаёт\n"
            f"- Что его интересует\n"
            f"- Как лучше с ним общаться"
        )
        ai_desc, _ = await get_best_ai_response(
            style_prompt, "Ты психолог и лингвист. Анализируй точно."
        )

        return {
            "analyzed": True,
            "tone": tone,
            "preferred_length": preferred_length,
            "avg_msg_len": round(avg_len),
            "total_messages": len(all_messages),
            "short_ratio": round(short_msgs / len(all_messages) * 100),
            "long_ratio": round(long_msgs / len(all_messages) * 100),
            "top_topics": all_topics[:8],
            "ai_description": ai_desc or "",
            "updated": datetime.now().strftime("%d.%m.%Y %H:%M")
        }
    except Exception as e:
        logger.error(f"analyze_style error: {e}")
        return {"analyzed": False, "reason": str(e)}


def build_personalized_system(data: dict) -> str:
    """Строит персонализированный системный промпт"""
    style = data.get("user_style", {})
    if not style.get("analyzed"):
        return ""

    parts = []
    tone = style.get("tone", "")
    length = style.get("preferred_length", "")
    desc = style.get("ai_description", "")

    if tone == "неформальный":
        parts.append("Общайся неформально, по-дружески, можно использовать простые слова.")
    elif tone == "формальный":
        parts.append("Общайся формально и профессионально.")

    if length == "краткий":
        parts.append("Давай краткие ответы — пользователь предпочитает короткие сообщения.")
    elif length == "подробный":
        parts.append("Давай подробные развёрнутые ответы — пользователь любит детали.")

    if desc:
        parts.append(f"Характер пользователя: {desc[:200]}")

    return "\n".join(parts)


# ============================================================
# ПОИСК МЕСТ РЯДОМ — геолокация + Nominatim/Overpass
# ============================================================

USER_LOCATIONS = {}  # chat_id -> {lat, lon, updated}

async def handle_location(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Получает геолокацию и предлагает что найти рядом"""
    loc = update.message.location
    chat_id = update.effective_chat.id
    USER_LOCATIONS[chat_id] = {
        "lat": loc.latitude,
        "lon": loc.longitude,
        "updated": datetime.now().strftime("%H:%M")
    }

    # Определяем адрес по координатам
    try:
        async with httpx.AsyncClient() as client:
            r = await client.get(
                f"https://nominatim.openstreetmap.org/reverse?lat={loc.latitude}&lon={loc.longitude}&format=json&accept-language=ru",
                headers={"User-Agent": "AIAgentBot/1.0"}, timeout=8
            )
            addr_data = r.json()
        addr = addr_data.get("display_name", "")
        city = addr_data.get("address", {}).get("city") or addr_data.get("address", {}).get("town", "")
        if city:
            data = load_data()
            data["weather_city"] = city
            save_data(data)
    except:
        addr = f"{loc.latitude:.4f}, {loc.longitude:.4f}"

    keyboard = InlineKeyboardMarkup([
        [InlineKeyboardButton("🍽 Кафе и рестораны", callback_data="nearby_food"),
         InlineKeyboardButton("☕ Кофейни", callback_data="nearby_cafe")],
        [InlineKeyboardButton("🏛 Музеи", callback_data="nearby_museum"),
         InlineKeyboardButton("🛒 Магазины", callback_data="nearby_shop")],
        [InlineKeyboardButton("💊 Аптеки", callback_data="nearby_pharmacy"),
         InlineKeyboardButton("🏦 Банки/ATM", callback_data="nearby_bank")],
        [InlineKeyboardButton("⛽ Заправки", callback_data="nearby_fuel"),
         InlineKeyboardButton("🏨 Отели", callback_data="nearby_hotel")],
        [InlineKeyboardButton("🔍 Свой запрос", callback_data="nearby_custom")],
    ])
    short_addr = addr[:80] + "..." if len(addr) > 80 else addr
    await update.message.reply_text(
        f"📍 *Геолокация получена!*\n\n{short_addr}\n\nЧто ищем рядом?",
        parse_mode="Markdown",
        reply_markup=keyboard
    )


async def search_nearby(lat: float, lon: float, category: str, label: str) -> str:
    """Ищет места рядом через Overpass API (OpenStreetMap)"""
    try:
        # Overpass запрос — радиус 1.5км
        radius = 1500
        overpass_queries = {
            "food": f'node["amenity"~"restaurant|cafe|fast_food|food_court"](around:{radius},{lat},{lon});',
            "cafe": f'node["amenity"="cafe"](around:{radius},{lat},{lon});',
            "museum": f'node["tourism"~"museum|gallery|attraction"](around:{radius},{lat},{lon});',
            "shop": f'node["shop"~"supermarket|mall|department_store|convenience"](around:{radius},{lat},{lon});',
            "pharmacy": f'node["amenity"="pharmacy"](around:{radius},{lat},{lon});',
            "bank": f'node["amenity"~"bank|atm"](around:{radius},{lat},{lon});',
            "fuel": f'node["amenity"="fuel"](around:{radius},{lat},{lon});',
            "hotel": f'node["tourism"~"hotel|hostel|motel"](around:{radius},{lat},{lon});',
        }
        query = overpass_queries.get(category, f'node["amenity"](around:{radius},{lat},{lon});')
        overpass_url = "https://overpass-api.de/api/interpreter"
        payload = f'[out:json][timeout:10];({query});out body 15;'

        async with httpx.AsyncClient() as client:
            r = await client.post(overpass_url, data=payload, timeout=12)
            data = r.json()

        elements = data.get("elements", [])
        if not elements:
            # Fallback через Nominatim
            return await search_nearby_nominatim(lat, lon, label)

        import math
        def dist(e):
            dlat = math.radians(e.get("lat",0) - lat)
            dlon = math.radians(e.get("lon",0) - lon)
            a = math.sin(dlat/2)**2 + math.cos(math.radians(lat))*math.cos(math.radians(e.get("lat",0)))*math.sin(dlon/2)**2
            return 6371000 * 2 * math.asin(math.sqrt(a))

        # Сортируем по расстоянию
        elements.sort(key=dist)
        lines = [f"📍 *{label} рядом с тобой:*\n"]

        for i, e in enumerate(elements[:7], 1):
            tags = e.get("tags", {})
            name = tags.get("name") or tags.get("name:ru") or tags.get("brand") or "Без названия"
            d_m = int(dist(e))
            d_str = f"{d_m}м" if d_m < 1000 else f"{d_m/1000:.1f}км"
            opening = tags.get("opening_hours", "")
            phone = tags.get("phone") or tags.get("contact:phone", "")
            rating = tags.get("stars", "")

            e_lat, e_lon = e.get("lat", lat), e.get("lon", lon)
            maps_url = f"https://www.openstreetmap.org/?mlat={e_lat}&mlon={e_lon}#map=18/{e_lat}/{e_lon}"
            route_url = f"https://www.openstreetmap.org/directions?from={lat},{lon}&to={e_lat},{e_lon}"

            line = f"*{i}. {name}* — {d_str}"
            if opening: line += f"\n   ⏰ {opening[:40]}"
            if phone: line += f"\n   📞 {phone}"
            line += f"\n   [📍 На карте]({maps_url}) | [🗺 Маршрут]({route_url})"
            lines.append(line)

        return "\n\n".join(lines)

    except Exception as e:
        logger.error(f"search_nearby error: {e}")
        return await search_nearby_nominatim(lat, lon, label)


async def search_nearby_nominatim(lat: float, lon: float, query: str) -> str:
    """Fallback поиск через Nominatim"""
    try:
        import urllib.parse
        url = f"https://nominatim.openstreetmap.org/search?q={urllib.parse.quote(query)}&format=json&limit=7&lat={lat}&lon={lon}&radius=2000&accept-language=ru"
        async with httpx.AsyncClient() as client:
            r = await client.get(url, headers={"User-Agent": "AIAgentBot/1.0"}, timeout=8)
            results = r.json()

        if not results:
            return f"❌ Рядом не найдено: {query}"

        import math
        lines = [f"📍 *{query} рядом:*\n"]
        for i, p in enumerate(results[:6], 1):
            p_lat, p_lon = float(p["lat"]), float(p["lon"])
            dlat = math.radians(p_lat - lat)
            dlon = math.radians(p_lon - lon)
            a = math.sin(dlat/2)**2 + math.cos(math.radians(lat))*math.cos(math.radians(p_lat))*math.sin(dlon/2)**2
            d_m = int(6371000 * 2 * math.asin(math.sqrt(a)))
            d_str = f"{d_m}м" if d_m < 1000 else f"{d_m/1000:.1f}км"
            name = p.get("name") or p["display_name"].split(",")[0]
            route_url = f"https://www.openstreetmap.org/directions?from={lat},{lon}&to={p_lat},{p_lon}"
            lines.append(f"*{i}. {name}* — {d_str}\n   [🗺 Маршрут]({route_url})")

        return "\n\n".join(lines)
    except Exception as e:
        return f"❌ Ошибка поиска: {e}"


NEARBY_CATEGORIES = {
    "nearby_food": ("food", "🍽 Кафе и рестораны"),
    "nearby_cafe": ("cafe", "☕ Кофейни"),
    "nearby_museum": ("museum", "🏛 Музеи"),
    "nearby_shop": ("shop", "🛒 Магазины"),
    "nearby_pharmacy": ("pharmacy", "💊 Аптеки"),
    "nearby_bank": ("bank", "🏦 Банки"),
    "nearby_fuel": ("fuel", "⛽ Заправки"),
    "nearby_hotel": ("hotel", "🏨 Отели"),
}



async def timezone_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Установить часовой пояс для утренней рассылки"""
    data = load_data()
    cur = data.get("timezone_offset", 5)
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("🇺🇿 UTC+5 Ташкент", callback_data="tz_5"),
         InlineKeyboardButton("🇷🇺 UTC+3 Москва", callback_data="tz_3")],
        [InlineKeyboardButton("🇰🇿 UTC+6 Алматы", callback_data="tz_6"),
         InlineKeyboardButton("🇦🇿 UTC+4 Баку", callback_data="tz_4")],
        [InlineKeyboardButton("🇺🇦 UTC+2 Киев", callback_data="tz_2"),
         InlineKeyboardButton("🇩🇪 UTC+1 Берлин", callback_data="tz_1")],
    ])
    await update.message.reply_text(
        f"🕐 *Часовой пояс*\n\nТекущий: UTC+{cur}\n\nВыбери свой город — бот будет присылать утреннюю сводку ровно в 7:00 по твоему времени:",
        parse_mode="Markdown", reply_markup=kb
    )

async def post_init(application: Application) -> None:
    """Действия после инициализации бота"""
    # Устанавливаем команды меню
    commands = [
        BotCommand("start", "🚀 Главное меню"),
        BotCommand("restart", "🔄 Перезапустить бота"),
        BotCommand("memory", "🧠 Память и история"),
        BotCommand("clear", "🗑 Очистить чат"),
        BotCommand("new", "💬 Новый чат"),
        BotCommand("chats", "📂 Список чатов"),
        BotCommand("search", "🌐 Поиск в интернете"),
        BotCommand("reminders", "⏰ Напоминания"),
        BotCommand("help", "❓ Помощь"),
        BotCommand("weather", "🌤 Погода"),
        BotCommand("summarize", "🔗 Пересказ ссылки"),
        BotCommand("tasks", "✅ Мои задачи"),
        BotCommand("stats", "📊 Статистика"),
        BotCommand("profile", "👤 Профиль"),
        BotCommand("goals", "🎯 Мои цели"),
        BotCommand("brain", "📚 База знаний"),
        BotCommand("tutor", "🎓 Репетитор"),
        BotCommand("habits", "🏃 Трекер привычек"),
        BotCommand("debate", "🤝 Дебаты"),
        BotCommand("interview", "💼 Собеседование"),
        BotCommand("books", "📖 Книжный клуб"),
        BotCommand("monitor", "🌐 Мониторинг сайтов"),
        BotCommand("export", "📤 Экспорт в PDF"),
        BotCommand("agent", "🤖 Агентный режим"),
        BotCommand("email", "📧 Email интеграция"),
        BotCommand("style", "🧠 Персонализация стиля"),
    ]
    
    await application.bot.set_my_commands(commands)
    
    # Запускаем фоновый цикл проверки напоминаний
    asyncio.create_task(reminder_check_loop(application))
    asyncio.create_task(morning_weather_job(application))
    asyncio.create_task(check_monitors_loop(application))

    logger.info("Бот успешно запущен!")

def main() -> None:
    """Главная функция запуска бота"""
    
    # Создаем приложение
    application = Application.builder()\
        .token(TELEGRAM_TOKEN)\
        .post_init(post_init)\
        .build()
    
    # Регистрируем обработчики команд
    application.add_handler(CommandHandler("start", start_command))
    application.add_handler(CommandHandler("restart", restart_command))
    application.add_handler(CommandHandler("memory", memory_command))
    application.add_handler(CommandHandler("clear", clear_command))
    application.add_handler(CommandHandler("new", new_chat_command))
    application.add_handler(CommandHandler("chats", chats_command))
    application.add_handler(CommandHandler("search", search_command))
    application.add_handler(CommandHandler("reminders", reminders_command))
    application.add_handler(CommandHandler("help", help_command))
    
    # Регистрируем обработчик инлайн кнопок
    application.add_handler(CallbackQueryHandler(button_callback))
    
    # Регистрируем обработчики сообщений
    application.add_handler(MessageHandler(filters.VOICE, handle_voice))
    application.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    application.add_handler(MessageHandler(filters.LOCATION, handle_location))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    application.add_handler(CommandHandler("weather", weather_command))
    application.add_handler(CommandHandler("summarize", summarize_command))
    application.add_handler(CommandHandler("tasks", tasks_command))
    application.add_handler(CommandHandler("stats", stats_command))
    application.add_handler(CommandHandler("profile", profile_command))
    application.add_handler(CommandHandler("goals", goals_command))
    application.add_handler(CommandHandler("brain", brain_command))
    application.add_handler(CommandHandler("tutor", tutor_command))
    application.add_handler(CommandHandler("habits", habits_command))
    application.add_handler(CommandHandler("debate", debate_command))
    application.add_handler(CommandHandler("interview", interview_command))
    application.add_handler(CommandHandler("books", books_command))
    application.add_handler(CommandHandler("monitor", monitor_command))
    application.add_handler(CommandHandler("export", export_command))
    application.add_handler(CommandHandler("agent", agent_command))
    application.add_handler(CommandHandler("email", email_command))
    application.add_handler(CommandHandler("style", finetune_command))
    application.add_handler(CommandHandler("timezone", timezone_command))
    
    # Запускаем бота
    logger.info("Запуск бота...")
    application.run_polling(
        drop_pending_updates=True,
        allowed_updates=["message", "callback_query"],
        close_loop=False
    )

if __name__ == "__main__":
    main()
